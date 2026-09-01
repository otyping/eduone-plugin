# -*- coding: utf-8 -*-
"""
build_slides.py — สร้าง .pptx จาก {BASE}_slides_<src>.json (EDU ONE — agent 4)

หลักการ
  * ใช้เทมเพลตของผู้ใช้เป็นฐาน (แยกตามวิชา × ชั้นปี) — กรอบตกแต่ง/ธีมมาจาก layout เดิม
  * หัวข้อ = Prompt · เนื้อหา/ตาราง = TH Sarabun New · ขนาดใหญ่พออ่านจากหลังห้อง
  * เนื้อหาไม่พอ 1 หน้า -> **ขึ้นหน้าใหม่ (ต่อ) โดยคงเนื้อหาเดิม** ห้ามย่อตัวอักษร
  * ตัดคำไทยตามขอบคำจริง (ZWSP) · สมการ $...$ เป็นยกกำลัง/ตัวห้อยจริงหรือ OMML
  * ฝังฟอนต์ลงไฟล์ท้ายสุดผ่าน PowerPoint COM (ปิดด้วย --no-embed)

ใช้:
  build_slides.py <slides.json> <out.pptx> [--template PATH] [--no-embed]

ใช้ Python 3.12: %LOCALAPPDATA%\\Programs\\Python\\Python312\\python.exe
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import media_cache as mc  # noqa: E402  (รูปจาก URL/พาธ -> ไฟล์ในเครื่อง)
import pptx_common as pc  # noqa: E402
import pptx_slots as ps  # noqa: E402  (เทมเพลตแบบสเปก LAYOUT)
import slides_media as sm  # noqa: E402  (แบบแผนชื่อไฟล์รูป + clean_prompt)

CONT = " (ต่อ)"
IMAGE_W_MM = 100          # กรอบภาพกว้าง 3.94 นิ้ว -> ย่อรูปที่ใหญ่เกินให้พอดีก่อนฝัง
BOTTOM_IN = 4.87          # ขอบล่างของพื้นที่เนื้อหา (กันชนกรอบตกแต่ง/เลขหน้า)
BASE_RE = re.compile(r"^([A-Z][0-9])-([A-Za-z]+)_")


def tokens_from_path(path):
    """ดึง GradeToken/SubjectToken จากชื่อไฟล์ BASE (เช่น P1-Sci_U1_1_slides_C1.pptx)"""
    m = BASE_RE.match(os.path.basename(path))
    return (m.group(1), m.group(2)) if m else ("", "")


# ---------------------------------------------------------------- ตัวช่วยแบ่งหน้า
def _pages_for(items, box, size=None):
    size = size or pc.SIZE["bullet"]
    return pc.fit_pages(items, pc.BODY_FONT, size, box[2], box[3])


def _body_box(top, width=None):
    return (pc.BODY_BOX[0], top, width or pc.BODY_BOX[2], BOTTOM_IN - top)


def _col_box(top):
    return (pc.COL_L_BOX[0], top, pc.COL_L_BOX[2], BOTTOM_IN - top)


def _image_box(top):
    return (pc.IMAGE_BOX[0], top, pc.IMAGE_BOX[2], BOTTOM_IN - top)


# ---------------------------------------------------------------- ทางเลือก A: layout หลากหลาย
# เปิดด้วย --variety : เปลี่ยนเฉพาะ "layout ที่ใช้เป็นพื้นหลัง" ตามบทบาทของหน้า
# กรอบข้อความ/ขนาดตัวอักษร/การแตกหน้ายังคำนวณเหมือนเดิมทุกอย่าง จึงไม่กระทบจำนวนหน้า
VARIETY = False
_VARIETY_KEY = {
    "objectives": "body_v1",
    "summary": "body_v2",
    "application": "body_v3",
}


def _key(base_key, section=""):
    """คืน layout key ที่จะใช้จริง — โหมดปกติใช้ค่าเดิม, โหมด variety สลับตามบทบาทของหน้า"""
    if not VARIETY:
        return base_key
    if base_key == "body":
        return _VARIETY_KEY.get(section, base_key)
    return base_key


def _content_top(title, wide=False):
    """y เริ่มเนื้อหา — เผื่อหัวข้อหน้าต่อ ' (ต่อ)' ที่อาจสูงกว่า เพื่อให้ทุกหน้าตรงกัน"""
    return max(pc.title_metrics(title, wide)[2],
               pc.title_metrics(title + CONT, wide)[2])


# ---------------------------------------------------------------- section builders
def build_cover(prs, s):
    slide = pc.add_slide(prs, "title")
    box = list(pc.COVER_TITLE_BOX)
    title = s.get("title", "")
    width_pt = (box[2] - 2 * pc.TEXT_INSET_IN) * 72
    size = pc.SIZE["cover_title"]
    while size > pc.SIZE["title"] and pc.wrap_count(title, pc.HEAD_FONT, size, width_pt, True) > 2:
        size -= 2
    shp = pc.add_textbox(slide, box, anchor=pc.MSO_ANCHOR.BOTTOM)
    par = shp.text_frame.paragraphs[0]
    pc._set_para(par, size, bullet=False, first=True, align=pc.PP_ALIGN.CENTER,
                 family=pc.HEAD_FONT)
    pc.add_text(par, title, pc.HEAD_FONT, size, bold=True, color=pc.INK)

    hook = s.get("hook_question", "")
    if hook:
        shp = pc.add_textbox(slide, pc.COVER_HOOK_BOX, anchor=pc.MSO_ANCHOR.TOP)
        par = shp.text_frame.paragraphs[0]
        pc._set_para(par, pc.SIZE["cover_hook"], bullet=False, first=True,
                     align=pc.PP_ALIGN.CENTER)
        pc.add_text(par, hook, pc.BODY_FONT, pc.SIZE["cover_hook"], color=pc.BODY_INK)
    return 1


def build_list(prs, s, items, bullet=True):
    """หน้าเนื้อหาแบบรายการเต็มความกว้าง (objectives / summary / application)"""
    title = s.get("title", "")
    lay = _key("body", s.get("section", ""))
    top = _content_top(title)
    box = _body_box(top)
    pages = _pages_for(items, box)
    for i, page in enumerate(pages):
        slide = pc.add_slide(prs, lay)
        pc.add_title(slide, title if i == 0 else title + CONT)
        pc.fill_bullets(pc.add_textbox(slide, box), page, bullet=bullet)
    return len(pages)


def _resolve_image(s, where):
    """คืน (พาธรูปในเครื่องที่ฝังได้, ค่า image_file ที่ขอมา) — ใช้ร่วมกันทุก section ที่มีรูป

    image_file เป็นได้ทั้งพาธสัมพัทธ์กับไฟล์ json และ URL — โหลดไม่ได้คืน "" แล้ววาดกรอบว่างแทน
    (ไม่ทำให้ build พัง)
    """
    want = s.get("image_file") or ""
    if not want:
        return "", ""
    cache = os.path.join(s.get("_media_dir", s.get("_dir", "")), ".cache")
    local = mc.resolve_image(want, IMAGE_W_MM, cache, base_dir=s.get("_dir", "")) or ""
    if not local:
        print("WARN: ใช้รูป %r ของหน้า %r ไม่ได้ — วางกรอบภาพว่างไว้แทน"
              % (want, where), file=sys.stderr)
    return local, want


def _image_note(s, want_img, image_file):
    """บรรทัดในโน้ตผู้สอน: ต้องใส่ไฟล์รูปชื่ออะไร + prompt ชุดเดียวกับใบสั่งผลิต"""
    if not s.get("image_prompt") or image_file:
        return ""
    prompt, _stripped = sm.clean_prompt(s["image_prompt"])
    return ("ภาพประกอบ: ใส่ไฟล์ %s ลงโฟลเดอร์ %s แล้ว build ใหม่\n"
            "prompt (คัดลอกไปวาง Google Nano Banana):\n%s"
            % (sm.asset_name(s.get("_stem", ""), s.get("_no", 0)),
               os.path.basename(s.get("_media_dir", "")), prompt))


def build_content(prs, s):
    """หน้าเนื้อหาหลัก — วางกรอบภาพคู่กับข้อความเมื่อทำได้ ไม่งั้นใช้เต็มความกว้าง

    กติกา (คงเนื้อหาเดิมเสมอ ไม่ย่อตัวอักษร):
      * ข้อความที่ใส่คอลัมน์แคบได้ >= 1 ข้อ -> หน้าแรกเป็น 2 คอลัมน์ (ข้อความ | ภาพ)
        ที่เหลือขึ้นหน้าใหม่เต็มความกว้าง
      * ข้อความยาวจนคอลัมน์แคบรับไม่ไหว -> ทุกหน้าเต็มความกว้าง แล้วเอากรอบภาพ
        ไปวางในที่ว่างท้ายหน้าสุดท้าย (ถ้าเหลือที่พอ) มิฉะนั้นเก็บ image_prompt ไว้ในโน้ต
    """
    title = s.get("title", "")
    bullets = s.get("bullets", [])
    image_file, want_img = _resolve_image(s, title)
    has_img = bool(s.get("image_prompt") or image_file)

    top = _content_top(title)
    col = _col_box(top)
    note_img = _image_note(s, want_img, image_file)

    side_by_side = []
    if has_img:
        cand = pc.fit_pages(bullets, pc.BODY_FONT, pc.SIZE["bullet"], col[2], col[3])[0]
        if pc.fits(cand, pc.BODY_FONT, pc.SIZE["bullet"], col[2], col[3]):
            side_by_side = cand

    made = 0
    if side_by_side:
        slide = pc.add_slide(prs, "two_col")
        pc.add_title(slide, title, wide=True)
        pc.fill_bullets(pc.add_textbox(slide, col), side_by_side)
        pc.add_image_frame(slide, _image_box(top), image_file=image_file or None)
        pc.set_notes(slide, s.get("speaker_note", ""), note_img)
        made = 1
    rest = bullets[len(side_by_side):]

    box = _body_box(top)
    pages = _pages_for(rest, box) if rest else []
    for k, page in enumerate(pages):
        slide = pc.add_slide(prs, "body")
        pc.add_title(slide, title + CONT if made else title)
        pc.fill_bullets(pc.add_textbox(slide, box), page)
        if not made:
            pc.set_notes(slide, s.get("speaker_note", ""), note_img if not has_img else "")
        made += 1
        # หน้าสุดท้ายเหลือที่ว่าง -> วางกรอบภาพลงไป
        if has_img and not side_by_side and k == len(pages) - 1:
            used_in = pc.block_height_pt(page, pc.BODY_FONT, pc.SIZE["bullet"], box[2]) / 72
            free = box[3] - used_in - 0.15
            if free >= 1.30:
                pc.add_image_frame(slide, (box[0], box[1] + box[3] - free, box[2], free),
                                   image_file=image_file or None)
            else:
                print("WARN: %r ไม่มีที่ว่างพอสำหรับกรอบภาพ — เก็บ image_prompt ไว้ในโน้ต"
                      % title, file=sys.stderr)
            pc.set_notes(slide, s.get("speaker_note", ""), note_img)
    return made


def build_question(prs, s):
    """หน้าคำถามฉายหน้าห้อง — ตัวคำถามอยู่บนจอจริง ไม่ใช่ประโยคบอกว่าครูจะถาม

    ใช้กับ L1/L2 เป็นหลัก (ขั้นนำ/คำถามต่อยอด) — มีภาพก็วางคู่กันสองคอลัมน์
    """
    text = (s.get("question") or "").strip()
    if not text:
        return 0
    title = (s.get("title") or "").strip()
    image_file, want_img = _resolve_image(s, title or text)
    has_img = bool(s.get("image_prompt") or image_file)

    slide = pc.add_slide(prs, "two_col" if has_img else _key("big_point"))
    top = pc.add_title(slide, title, wide=has_img) if title else pc.BODY_BOX[1]
    box = (_col_box(top) if has_img else _body_box(top))

    # เลือกขนาดที่ใหญ่ที่สุดที่ยังอยู่ในกรอบ (ไม่แตกหน้า — คำถามต้องอยู่หน้าเดียว)
    size = pc.SIZE["question"] if has_img else pc.SIZE["question_solo"]
    while size > pc.SIZE["question_min"] and not pc.fits(
            [text], pc.BODY_FONT, size, box[2], box[3], bullet=False):
        size -= 2
    shp = pc.add_textbox(slide, box, anchor=pc.MSO_ANCHOR.MIDDLE)
    pc.fill_bullets(shp, [text], size_pt=size, bullet=False, color=pc.INK,
                    align=pc.PP_ALIGN.CENTER if not has_img else pc.PP_ALIGN.LEFT)
    if has_img:
        pc.add_image_frame(slide, _image_box(top), image_file=image_file or None)
    pc.set_notes(slide, s.get("speaker_note", ""), _image_note(s, want_img, image_file))
    return 1


def build_vocab(prs, s):
    """ตารางคำศัพท์ — หัวตารางซ้ำทุกหน้า, สูงสุด 6 คำต่อหน้า ตาม Master Prompt"""
    table = s.get("table", [])
    if not table:
        return 0
    head, rows = table[0], table[1:]
    title = s.get("title", "คำศัพท์ภาษาอังกฤษน่ารู้")
    top = _content_top(title, wide=True)
    box = (pc.TABLE_BOX[0], top, pc.TABLE_BOX[2], min(pc.TABLE_BOX[3], BOTTOM_IN - top))
    per_page = min(6, pc.table_rows_fit(table, box[3]) - 1)
    chunks = [rows[i:i + per_page] for i in range(0, len(rows), per_page)] or [[]]
    row_h_in = pc.table_row_h_in()
    for i, chunk in enumerate(chunks):
        slide = pc.add_slide(prs, "title_only")
        pc.add_title(slide, title if i == 0 else title + CONT, wide=True)
        h = min(box[3], (len(chunk) + 1) * row_h_in)
        pc.add_table(slide, [head] + chunk, (box[0], box[1], box[2], h))
    return len(chunks)


BUILDERS = {
    "cover": lambda prs, s: build_cover(prs, s),
    "objectives": lambda prs, s: build_list(prs, s, s.get("items", []), bullet=False),
    "content": lambda prs, s: build_content(prs, s),
    "question": lambda prs, s: build_question(prs, s),
    "summary": lambda prs, s: build_list(prs, s, s.get("bullets", [])),
    "application": lambda prs, s: build_list(prs, s, s.get("bullets", [])),
    "vocab": lambda prs, s: build_vocab(prs, s),
}


# ---------------------------------------------------------------- main
def build(spec_path, out_path, template=None, embed=True):
    with open(spec_path, "r", encoding="utf-8-sig") as f:
        spec = json.load(f)

    grade, subject = tokens_from_path(out_path)
    if not grade:
        # ชื่อไฟล์ปลายทางอาจไม่ได้ตั้งตาม BASE (เช่นตอน build ลงไฟล์ชั่วคราวเพื่อนับหน้า)
        # ถ้าไม่ถอยมาอ่านจากชื่อ spec ต่อ resolver จะเงียบ ๆ ไปหยิบเทมเพลตของวิชาอื่นมาใช้
        # แล้วให้จำนวนหน้า/หน้าตาผิดโดยไม่มีใครรู้
        grade, subject = tokens_from_path(spec_path)
    if not template:
        template = pc.resolve_template(subject, grade)
    print("เทมเพลต: %s" % template)
    # ปรับตัวตามเทมเพลตก่อนวาด: จานสีหมึกตามความสว่างของพื้น + เลือก layout ที่ลวดลาย
    # ตกแต่งไม่ทับพื้นที่ข้อความ (วัดจากภาพจริงครั้งเดียวต่อเทมเพลตแล้วแคชไว้)
    pc.use_template(template)

    spec_dir = os.path.dirname(os.path.abspath(spec_path))
    stem = sm.stem_of(spec_path)
    media_dir = sm.media_dir_of(spec_path)
    for i, s in enumerate(spec.get("slides", []), 1):
        s["_dir"] = spec_dir
        s["_media_dir"] = media_dir
        s["_stem"] = stem
        s["_no"] = i

    # เทมเพลตมี 2 แบบ — เลือก engine ตามชนิดของเทมเพลต
    #   (ก) แบบสเปก LAYOUT: หน้าตัวอย่าง 1 หน้าต่อ 1 LAYOUT + วัตถุตั้งชื่อ -> โคลนแล้วเติมช่อง
    #   (ข) แบบเดิม: ดีไซน์อยู่ใน slideLayout -> ลบสไลด์ตัวอย่างแล้ววาดกล่องเอง
    from pptx import Presentation as _Presentation
    raw = _Presentation(template)
    if ps.is_slot_template(raw):
        import build_slides_slots as bss
        print("engine: สเปก LAYOUT (%d layout ในเทมเพลต)" % len(ps.layout_index(raw)))
        total = bss.build(raw, spec, spec_path,
                          lambda sl: _resolve_image(sl, sl.get("title", "")),
                          lambda sl: _image_note(sl, sl.get("image_file", ""),
                                                 _resolve_image(sl, "")[0]))
        os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
        raw.save(out_path)
        print("เขียนแล้ว: %s (%d สไลด์)" % (out_path, total))
        if embed:
            import embed_fonts_pptx
            if embed_fonts_pptx.embed(out_path) != 0:
                print("WARN: ฝังฟอนต์ไม่สำเร็จ", file=sys.stderr)
        return total

    prs = pc.open_template(template)
    total = 0
    for s in spec.get("slides", []):
        fn = BUILDERS.get(s.get("section"))
        if fn is None:
            print("WARN: ข้าม section ที่ไม่รู้จัก: %r" % s.get("section"), file=sys.stderr)
            continue
        n = fn(prs, s)
        total += n
        print("  %-12s -> %d หน้า : %s" % (s.get("section"), n, s.get("title", "")[:40]))

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    print("เขียนแล้ว: %s (%d สไลด์)" % (out_path, total))

    if embed:
        import embed_fonts_pptx
        if embed_fonts_pptx.embed(out_path) != 0:
            print("WARN: ฝังฟอนต์ไม่สำเร็จ — ไฟล์ยังใช้ได้ แต่เครื่องปลายทางต้องมีฟอนต์เอง",
                  file=sys.stderr)
    return total


def main():
    ap = argparse.ArgumentParser(description="สร้าง .pptx จาก slides.json")
    ap.add_argument("spec")
    ap.add_argument("out")
    ap.add_argument("--template", help="ระบุเทมเพลตเอง (ปกติหาให้จากวิชา×ชั้น)")
    ap.add_argument("--no-embed", action="store_true", help="ไม่ต้องฝังฟอนต์")
    ap.add_argument("--variety", action="store_true",
                    help="ทางเลือก A: สลับ layout พื้นหลังตามบทบาทของหน้า (กรอบข้อความเหมือนเดิม)")
    ap.add_argument("--engine", choices=["auto", "boxes", "placeholders"], default="auto",
                    help="ทางเลือก B: 'placeholders' = วางเนื้อหาลง placeholder ของ layout จริง")
    a = ap.parse_args()
    global VARIETY
    VARIETY = a.variety
    if a.engine == "placeholders":
        import build_slides_ph as bph
        return bph.build(a.spec, a.out, a.template, embed=not a.no_embed)
    build(a.spec, a.out, a.template, embed=not a.no_embed)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
