# -*- coding: utf-8 -*-
"""
build_slides_ph.py — เอนจินสไลด์ "ทางเลือก B": วางเนื้อหาลง **placeholder จริงของ layout**

ต่างจากเอนจินหลัก (`build_slides.py`) อย่างไร
  * เอนจินหลัก: ลบ placeholder ทิ้งทั้งหมด แล้ววาดกล่องเองที่พิกัดคงที่ทุกหน้า
    -> คุมฟอนต์/ขนาด/การแตกหน้าได้แม่น แต่ใช้เทมเพลตเป็น "วอลเปเปอร์" เท่านั้น
    ทุกหน้าจึงมีองค์ประกอบเหมือนกันหมด ไม่ว่าเทมเพลตจะออกแบบ layout มากี่แบบ
  * เอนจินนี้: เลือก layout ตาม **รูปร่างของเนื้อหา** แล้วเขียนลงช่องที่ดีไซเนอร์วางไว้
    -> ได้องค์ประกอบตามที่เทมเพลตออกแบบมา (หัวข้อใหญ่กลางจอ · ช่องภาพด้านขวา ฯลฯ)

ยังคงกติกาเดิมทุกข้อ: ฟอนต์ Prompt/TH Sarabun New (ตั้ง cs ด้วย) · ตัดคำไทยด้วย ZWSP ·
ขนาดอ่านจากหลังห้อง · **ล้นกรอบ = ขึ้นหน้าใหม่ "(ต่อ)" ไม่ย่อตัวอักษร** (วัดจากกรอบของ placeholder)

ใช้ผ่าน: build_slides.py <spec> <out> --engine placeholders

ใช้ Python 3.12: %LOCALAPPDATA%\\Programs\\Python\\Python312\\python.exe
"""
from __future__ import annotations

import json
import os
import sys

from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pptx_common as pc  # noqa: E402

CONT = " (ต่อ)"

# บทบาทของหน้า -> ชื่อ layout ที่อยากได้ (ไล่ตามลำดับ ไม่เจอก็ถอยไปตัวถัดไป)
# ตัวสุดท้ายของทุกแถวต้องเป็น layout พื้นฐานที่เทมเพลตไหนก็มี
ROLE_LAYOUTS = {
    "cover":       ("TITLE", "TITLE_AND_BODY"),
    "objectives":  ("CUSTOM_4", "TITLE_AND_BODY"),          # หัวข้อ + ก้อนข้อความใหญ่
    "content":     ("TITLE_AND_BODY",),
    "content_img": ("CUSTOM_5", "TITLE_AND_BODY"),          # หัวข้อ + ข้อความ + ช่องภาพขวา
    "question":    ("MAIN_POINT", "TITLE_AND_BODY"),        # หัวข้อใหญ่กลางจอ
    "question_img": ("CUSTOM_5", "TITLE_AND_BODY"),
    "summary":     ("TITLE_AND_BODY",),
    "application": ("CUSTOM_4", "TITLE_AND_BODY"),
    "vocab":       ("TITLE_ONLY", "TITLE_AND_BODY"),
}

TITLE_TYPES = ("CENTER_TITLE", "TITLE")
BODY_TYPES = ("BODY", "SUBTITLE", "OBJECT")
PIC_TYPES = ("PICTURE",)


def _t(ph):
    try:
        return str(ph.placeholder_format.type).split(" ")[0]
    except Exception:
        return "?"


def _box(shape):
    """กรอบของรูปทรงเป็นนิ้ว (x, y, w, h)"""
    return ((shape.left or 0) / 914400.0, (shape.top or 0) / 914400.0,
            (shape.width or 0) / 914400.0, (shape.height or 0) / 914400.0)


def _set_box(shape, x, y, w, h):
    """ตั้งกรอบของ placeholder — **ต้องเขียนครบทั้ง 4 ค่าเสมอ**

    placeholder ที่สืบทอดตำแหน่ง/ขนาดมาจาก layout ยังไม่มี `a:xfrm` ของตัวเอง
    ถ้าสั่งแค่ค่าเดียว (เช่น .height) python-pptx จะสร้าง xfrm ที่มีแต่ค่านั้น
    ที่เหลือเป็น None -> PowerPoint ตีความเป็น 0 แล้วกล่องยุบจนอ่านไม่ได้
    (เคยเกิดจริง: กล่องหัวข้อเหลือกว้างเกือบศูนย์ ตัวอักษรไทยเรียงลงมาทีละตัว)
    """
    shape.left, shape.top = int(x * 914400), int(y * 914400)
    shape.width, shape.height = int(w * 914400), int(h * 914400)


def _pick_layout(prs, role):
    by_name = {L.name: L for L in prs.slide_masters[0].slide_layouts}
    for name in ROLE_LAYOUTS.get(role, ()):
        if name in by_name:
            return by_name[name], name
    return prs.slide_masters[0].slide_layouts[0], prs.slide_masters[0].slide_layouts[0].name


def _new(prs, role):
    """เพิ่มสไลด์จาก layout ของบทบาทนั้น **โดยคง placeholder ไว้** แล้วคืน (slide, ช่องที่จับได้)"""
    layout, name = _pick_layout(prs, role)
    sl = prs.slides.add_slide(layout)
    slots = {"title": None, "body": None, "pic": None, "extra": []}
    for ph in sl.placeholders:
        t = _t(ph)
        if t in TITLE_TYPES and slots["title"] is None:
            slots["title"] = ph
        elif t in BODY_TYPES and slots["body"] is None:
            slots["body"] = ph
        elif t in PIC_TYPES and slots["pic"] is None:
            slots["pic"] = ph
        else:
            slots["extra"].append(ph)
    return sl, slots, name


def _prep(ph, anchor=MSO_ANCHOR.TOP):
    """เตรียม placeholder ให้พฤติกรรมเหมือนกล่องข้อความของเรา (สำคัญ: ห้ามย่อตัวอักษรเอง)"""
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    bodyPr = tf._txBody.bodyPr
    for attr in ("lIns", "rIns", "tIns", "bIns"):
        bodyPr.set(attr, str(int(pc.TEXT_INSET_IN * 914400)))
    return ph


def _drop(shape):
    shape._element.getparent().remove(shape._element)


def _cleanup(sl, used):
    """ลบ placeholder ที่ไม่ได้ใช้ ไม่งั้นจะเหลือกล่อง 'Click to add text' ติดไปในไฟล์

    เก็บเป็น XML element เพราะวัตถุ SlidePlaceholder ของ python-pptx ไม่ hashable
    """
    for ph in list(sl.placeholders):
        if ph._element not in used:
            _drop(ph)


def _put_title(slots, text, used, size_key="title"):
    """วางหัวข้อลงช่องของ layout — ถ้ายังไม่พอจะ **ขยายช่องแล้วดันกรอบเนื้อหาลง**

    ช่องหัวข้อของเทมเพลตสูงคงที่ (มัก 0.6 นิ้ว = 1 บรรทัด) หัวข้อยาวจึงล้นกรอบ
    เอนจินหลักแก้ด้วยการขยายกล่องหัวข้อ ที่นี่ก็ต้องทำเหมือนกัน ไม่งั้น verify_pptx ตก
    """
    ph = slots["title"]
    if ph is None or not text:
        return
    x, y, w, h = _box(ph)
    width_pt = (w - 2 * pc.TEXT_INSET_IN) * 72
    size = pc.SIZE[size_key]
    floor = pc.SIZE["title_min"] if size_key == "title" else 24

    def need(sz):
        lines = pc.wrap_count(text, pc.HEAD_FONT, sz, width_pt, True)
        return lines * pc.line_h_pt(pc.HEAD_FONT, sz) / 72 + 2 * pc.TEXT_INSET_IN

    while size > floor and need(size) > h:
        size -= 2
    grow = need(size) - h
    if grow > 0.001:                      # เล็กสุดแล้วยังไม่พอ -> ขยายช่องหัวข้อลงมา
        _set_box(ph, x, y, w, h + grow)
        body = slots.get("body")
        if body is not None:
            bx, by, bw, bh = _box(body)
            shift = max(0.0, (y + h + grow + 0.10) - by)
            if shift > 0 and bh - shift > 0.6:
                _set_box(body, bx, by + shift, bw, bh - shift)
    _prep(ph, MSO_ANCHOR.MIDDLE if size_key != "title" else MSO_ANCHOR.TOP)
    par = ph.text_frame.paragraphs[0]
    pc._set_para(par, size, bullet=False, first=True, family=pc.HEAD_FONT,
                 align=PP_ALIGN.CENTER if size_key != "title" else PP_ALIGN.LEFT)
    pc.add_text(par, text, pc.HEAD_FONT, size, bold=True, color=pc.INK)
    used.add(ph._element)


def _body_pages(slots, items, size=None):
    """แบ่งหน้าโดยวัดกับกรอบของ placeholder จริง (ไม่ใช่กรอบคงที่)"""
    ph = slots["body"]
    if ph is None or not items:
        return [], None
    x, y, w, h = _box(ph)
    size = size or pc.SIZE["bullet"]
    return pc.fit_pages(items, pc.BODY_FONT, size, w, h), (x, y, w, h)


def _put_body(slots, page, used, size=None, bullet=True, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    ph = slots["body"]
    if ph is None or not page:
        return
    _prep(ph, anchor)
    pc.fill_bullets(ph, page, size_pt=size or pc.SIZE["bullet"], bullet=bullet,
                    color=pc.BODY_INK, align=align)
    used.add(ph._element)


def _put_image(sl, slots, used, image_file, caption="ภาพประกอบ"):
    """ใช้กรอบของ PICTURE placeholder เป็นที่วางรูป/กรอบภาพ แล้วลบ placeholder เดิมทิ้ง"""
    ph = slots["pic"]
    if ph is None:
        return False
    box = _box(ph)
    _drop(ph)                       # วาดกรอบของเราเองทับตำแหน่งเดิม (คุมเส้น/ป้ายได้)
    slots["pic"] = None
    pc.add_image_frame(sl, box, caption=caption, image_file=image_file or None)
    return True


# ================================================================ ตัวสร้างรายหน้า
def _cover(prs, s, ctx):
    sl, slots, lay = _new(prs, "cover")
    used = set()
    _put_title(slots, s.get("title", ""), used, size_key="cover_title")
    hook = s.get("hook_question", "")
    if hook and slots["body"] is not None:
        _put_body(slots, [hook], used, size=pc.SIZE["cover_hook"], bullet=False,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    _cleanup(sl, used)
    return 1, lay


def _list(prs, s, ctx, items, bullet=True, role=None):
    role = role or s.get("section", "content")
    title = s.get("title", "")
    made = 0
    lay = ""
    # แบ่งหน้าโดยดูกรอบจริงของ layout ที่จะใช้ (ต้องสร้างหน้าแรกก่อนถึงจะรู้กรอบ)
    sl, slots, lay = _new(prs, role)
    used = set()
    _put_title(slots, title, used)          # ต้องมาก่อน: อาจขยายหัวข้อแล้วย่อกรอบเนื้อหา
    pages, _b = _body_pages(slots, items)
    if not pages:
        pages = [[]]
    _put_body(slots, pages[0], used, bullet=bullet)
    _cleanup(sl, used)
    pc.set_notes(sl, s.get("speaker_note", ""))
    made += 1
    for page in pages[1:]:
        sl, slots, lay = _new(prs, role)
        used = set()
        _put_title(slots, title + CONT, used)
        _put_body(slots, page, used, bullet=bullet)
        _cleanup(sl, used)
        made += 1
    return made, lay


def _content(prs, s, ctx):
    title = s.get("title", "")
    bullets = s.get("bullets", [])
    image_file, want = ctx["resolve"](s)
    has_img = bool(s.get("image_prompt") or image_file)
    role = "content_img" if has_img else "content"

    # หน้าแรกอาจใช้ layout ที่มีช่องภาพ (ช่องข้อความแคบกว่ามาก) ส่วนหน้าถัดไปเป็น layout
    # เต็มความกว้าง -> **ต้องแบ่งหน้าที่เหลือใหม่กับกรอบของ layout นั้น** ไม่ใช่ยกผลจากกล่องแคบมาใช้ต่อ
    # (ถ้ายกมาใช้ต่อ เนื้อหาจะกระจายหน้าละข้อ เหลือที่ว่างครึ่งจอทุกหน้า)
    sl, slots, lay = _new(prs, role)
    used = set()
    _put_title(slots, title, used)          # ต้องมาก่อน: อาจขยายหัวข้อแล้วย่อกรอบเนื้อหา
    first, _b = _body_pages(slots, bullets)
    first = first[0] if first else []
    _put_body(slots, first, used)
    if has_img:
        _put_image(sl, slots, used, image_file)
    _cleanup(sl, used)
    pc.set_notes(sl, s.get("speaker_note", ""), ctx["note"](s, want, image_file))
    made, used_lays = 1, [lay]

    rest = bullets[len(first):]
    while rest:
        sl, slots, lay = _new(prs, "content")
        used = set()
        _put_title(slots, title + CONT, used)
        pages, _b = _body_pages(slots, rest)
        page = pages[0] if pages else rest
        _put_body(slots, page, used)
        _cleanup(sl, used)
        made += 1
        used_lays.append(lay)
        if len(page) >= len(rest):
            break
        rest = rest[len(page):]
    return made, "+".join(sorted(set(used_lays)))


def _question(prs, s, ctx):
    text = (s.get("question") or "").strip()
    if not text:
        return 0, ""
    title = (s.get("title") or "").strip()
    image_file, want = ctx["resolve"](s)
    has_img = bool(s.get("image_prompt") or image_file)
    sl, slots, lay = _new(prs, "question_img" if has_img else "question")
    used = set()
    if has_img:
        _put_title(slots, title or "ลองคิดดู", used)
        _put_body(slots, [text], used, size=pc.SIZE["question"], bullet=False)
        _put_image(sl, slots, used, image_file)
    else:
        # MAIN_POINT: ช่องหัวข้อคือกล่องใหญ่กลางจอ -> เอาตัวคำถามไปไว้ตรงนั้น
        _put_title(slots, text, used, size_key="question_solo")
        if title and slots["body"] is not None:
            _put_body(slots, [title], used, size=pc.SIZE["bullet"], bullet=False,
                      align=PP_ALIGN.CENTER)
    _cleanup(sl, used)
    pc.set_notes(sl, s.get("speaker_note", ""), ctx["note"](s, want, image_file))
    return 1, lay


def _vocab(prs, s, ctx):
    table = s.get("table", [])
    if not table:
        return 0, ""
    head, rows = table[0], table[1:]
    title = s.get("title", "คำศัพท์ภาษาอังกฤษน่ารู้")
    made, lay = 0, ""
    per_page = 6
    chunks = [rows[i:i + per_page] for i in range(0, len(rows), per_page)] or [[]]
    row_h = pc.table_row_h_in()
    for i, chunk in enumerate(chunks):
        sl, slots, lay = _new(prs, "vocab")
        used = set()
        _put_title(slots, title if i == 0 else title + CONT, used)
        # ตารางวางใต้หัวข้อ โดยอิงกรอบหัวข้อของ layout เอง
        tx, ty, tw, th = _box(slots["title"]) if slots["title"] is not None \
            else (0.79, 0.42, 8.43, 0.8)
        top = ty + th + 0.15
        h = min(4.87 - top, (len(chunk) + 1) * row_h)
        _cleanup(sl, used)
        pc.add_table(sl, [head] + chunk, (tx, top, tw, h))
        made += 1
    return made, lay


BUILDERS = {
    "cover": lambda prs, s, c: _cover(prs, s, c),
    "objectives": lambda prs, s, c: _list(prs, s, c, s.get("items", []), bullet=False),
    "content": lambda prs, s, c: _content(prs, s, c),
    "question": lambda prs, s, c: _question(prs, s, c),
    "summary": lambda prs, s, c: _list(prs, s, c, s.get("bullets", [])),
    "application": lambda prs, s, c: _list(prs, s, c, s.get("bullets", [])),
    "vocab": lambda prs, s, c: _vocab(prs, s, c),
}


def build(spec_path, out_path, template=None, embed=True):
    import build_slides as bs   # ใช้ตัวช่วยเรื่องรูป/เทมเพลตร่วมกัน ไม่เขียนซ้ำ

    with open(spec_path, encoding="utf-8-sig") as f:
        spec = json.load(f)

    grade, subject = bs.tokens_from_path(out_path)
    if not grade:
        grade, subject = bs.tokens_from_path(spec_path)
    if not template:
        template = pc.resolve_template(subject, grade)
    print("เทมเพลต: %s" % template)
    print("engine: placeholders (ทางเลือก B — วางเนื้อหาลงช่องของ layout)")
    pc.use_template(template)

    spec_dir = os.path.dirname(os.path.abspath(spec_path))
    import slides_media as sm
    stem, media_dir = sm.stem_of(spec_path), sm.media_dir_of(spec_path)
    for i, s in enumerate(spec.get("slides", []), 1):
        s["_dir"], s["_media_dir"], s["_stem"], s["_no"] = spec_dir, media_dir, stem, i

    ctx = {"resolve": lambda s: bs._resolve_image(s, s.get("title", "")),
           "note": lambda s, w, f: bs._image_note(s, w, f)}

    prs = pc.open_template(template)
    total, seen = 0, {}
    for s in spec.get("slides", []):
        fn = BUILDERS.get(s.get("section"))
        if fn is None:
            print("WARN: ข้าม section ที่ไม่รู้จัก: %r" % s.get("section"), file=sys.stderr)
            continue
        n, lay = fn(prs, s, ctx)
        total += n
        seen[lay] = seen.get(lay, 0) + n
        print("  %-12s -> %d หน้า [%s] : %s"
              % (s.get("section"), n, lay, s.get("title", "")[:34]))

    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    prs.save(out_path)
    print("เขียนแล้ว: %s (%d สไลด์)" % (out_path, total))
    print("ใช้ layout: %s" % ", ".join("%s x%d" % (k, v) for k, v in sorted(seen.items())))

    if embed:
        import embed_fonts_pptx
        if embed_fonts_pptx.embed(out_path) != 0:
            print("WARN: ฝังฟอนต์ไม่สำเร็จ — ไฟล์ยังใช้ได้ แต่เครื่องปลายทางต้องมีฟอนต์เอง",
                  file=sys.stderr)
    return total
