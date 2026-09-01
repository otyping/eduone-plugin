# -*- coding: utf-8 -*-
"""
pptx_common.py — เลเยอร์กลางของสไลด์ .pptx (EDU ONE) — คู่ขนานกับ docx_common.py

หน้าที่
  1. หาเทมเพลตตาม (วิชา × ชั้น) แล้วเปิดเป็นฐาน (ล้างสไลด์ตัวอย่างออกให้หมด)
  2. ตั้งฟอนต์ให้ถูกทั้ง latin **และ cs (complex script = ไทย)** — ถ้าไม่ตั้ง cs
     ข้อความไทยจะตกไปใช้ฟอนต์ธีมเวลาเปิดเครื่องอื่น
  3. ตัดคำไทยด้วย ZWSP (pythainlp newmm) เหมือนที่ .docx ใช้
  4. วัดความกว้าง/ความสูงข้อความจากไฟล์ฟอนต์จริง (PIL + ตาราง hhea)
     -> รู้ล่วงหน้าว่าข้อความล้นกรอบไหม เพื่อ **ขึ้นหน้าใหม่ ไม่ใช่ย่อตัวอักษร**
  5. สัญลักษณ์คณิต/วิทย์: $...$ -> run ยกกำลัง/ตัวห้อยจริง (หรือ OMML เมื่อมี \\frac \\sqrt)

รันผ่านตัวห่อ: eduone-py <ชื่อไฟล์นี้> <args>  (หา Python 3.12 ให้เองทุก OS)
"""
from __future__ import annotations

import copy
import functools
import glob
import json
import os
import re
import struct
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import docx_common as dc  # noqa: E402  (ใช้ pipeline อินไลน์/ตัดคำร่วมกับ .docx)

# ---------------------------------------------------------------- ฟอนต์
HEAD_FONT = "Prompt"            # หัวข้อทุกชนิด
BODY_FONT = "TH Sarabun New"    # เนื้อหา/ตาราง
FONT_FAMILIES = (HEAD_FONT, BODY_FONT)

_FONT_FILES = {
    (HEAD_FONT, False): "Prompt-Regular.ttf",
    (HEAD_FONT, True): "Prompt-Bold.ttf",
    (BODY_FONT, False): "THSarabunNew.ttf",
    (BODY_FONT, True): "THSarabunNew Bold.ttf",
}
_FONT_DIRS = [
    os.path.expandvars(r"%LOCALAPPDATA%\Microsoft\Windows\Fonts"),
    os.path.expandvars(r"%WINDIR%\Fonts"),
]

# ---------------------------------------------------------------- ขนาด (โปรไฟล์ "ใหญ่")
# สไลด์เทมเพลตกว้าง 10 นิ้ว -> ทุกค่าเทียบเท่า x1.33 บนสไลด์มาตรฐาน 13.3 นิ้ว
SIZE = {
    "cover_title": 40,   # เทียบเท่า 53pt
    "cover_hook": 26,    # 35pt
    "title": 32,         # 43pt
    "title_min": 26,     # หัวข้อยาวลดได้ถึงเท่านี้ (ยังเทียบเท่า 35pt)
    "bullet": 28,        # 37pt
    "bullet_sub": 24,    # 32pt
    "table_head": 24,
    "table_cell": 24,
    "caption": 24,
    "question": 34,      # คำถามฉายหน้าห้อง (L1/L2) คู่กับภาพ — ใหญ่กว่าบุลเล็ต
    "question_solo": 40, # คำถามเดี่ยวเต็มจอ (ไม่มีภาพ) — มีที่ว่างทั้งหน้า ใหญ่ได้เต็มที่
    "question_min": 28,
    "min_body": 24,      # พื้นขนาดต่ำสุดของเนื้อหา — ต่ำกว่านี้ให้ขึ้นหน้าใหม่แทน
    "min_table": 22,
}

# ---------------------------------------------------------------- สี (จากพาเลตเทมเพลต)
INK = RGBColor(0x0E, 0x2A, 0x47)       # น้ำเงินเข้ม — หัวข้อ
BODY_INK = RGBColor(0x20, 0x20, 0x20)  # เกือบดำ — เนื้อหา
# หมายเหตุ: สี่ตัวนี้ (INK/BODY_INK/SOFT/MUTED) เป็นค่า "เริ่มต้นสำหรับเทมเพลตพื้นสว่าง"
# และจะถูกสลับเป็นจานสีพื้นเข้มอัตโนมัติโดย apply_palette() — ดู detect_bg_dark() ท้ายไฟล์
BRAND = RGBColor(0x64, 0x6B, 0x96)     # #646b96 หัวตาราง
BAND = RGBColor(0xF3, 0xF3, 0xF3)      # #f3f3f3 แถบตาราง
# หมึกสำหรับข้อความที่วางบน "พื้นสีอ่อนที่เราวาดเอง" (เซลล์ตาราง · กรอบภาพ)
# ต้องเข้มเสมอ ไม่เปลี่ยนตามจานสีของพื้นสไลด์ ไม่งั้นบนเทมเพลตพื้นเข้มจะกลายเป็นอ่อนบนอ่อน
TABLE_INK = RGBColor(0x20, 0x20, 0x20)
FRAME_MUTED = RGBColor(0x86, 0x9F, 0xB2)
SOFT = RGBColor(0xD5, 0xD9, 0xEE)      # #d5d9ee เส้นกรอบภาพ
MUTED = RGBColor(0x86, 0x9F, 0xB2)     # ข้อความจาง (ป้ายกรอบภาพ)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------- เรขาคณิต (นิ้ว)
# ระยะบรรทัดแบบ "แน่นอน" (a:lnSpc spcPts) — ไม่ใช้ % เพราะ PowerPoint จะยืดบรรทัด
# ที่มีสระ/วรรณยุกต์ซ้อนหรือวงเล็บละตินสูงถึง 1.57 เท่า ทำให้กะหน้าไม่ได้และจังหวะบรรทัดกระโดด
# ค่าที่ตั้งมาจากการวัด ink extent จริงของฟอนต์ (ไทยซ้อนสูงสุด 1.05 em ทั้งสองฟอนต์)
LINE_EXACT = {"TH Sarabun New": 1.18, "Prompt": 1.20}
SPC_BEF_RATIO = 0.35     # ระยะก่อนย่อหน้า = 0.35 x ขนาดตัวอักษร
TEXT_INSET_IN = 0.05     # ระยะขอบในกล่องข้อความ (บน+ล่าง / ซ้าย+ขวา อย่างละด้าน)
BULLET_INDENT_IN = 0.30  # ระยะย่อหน้าแขวนของบุลเล็ต
# วัดเทียบกับที่ PowerPoint จัดจริงแล้ว (PIL 481.9pt vs render 479.7pt = คลาด 0.5%)
# ระยะบรรทัดเป็นค่าคงที่ (spcPts) ความสูงจึงคำนวณตรง เหลือเผื่อไว้กันปัดเศษเท่านั้น
WIDTH_SAFETY = 0.98
HEIGHT_SAFETY = 1.00

TITLE_BOX = (0.79, 0.42, 7.66, 0.80)          # x, y, w, h — หัวข้อหน้า (คอลัมน์เดียว)
TITLE_BOX_WIDE = (0.79, 0.42, 8.43, 0.80)     # หัวข้อหน้าแบบกว้าง (ตาราง/2 คอลัมน์)
BODY_BOX = (0.79, 1.32, 7.66, 3.55)           # เนื้อหาเต็มความกว้าง
COL_L_BOX = (0.79, 1.32, 4.25, 3.55)          # เนื้อหาคอลัมน์ซ้าย (คู่กับกรอบภาพ)
IMAGE_BOX = (5.28, 1.32, 3.94, 3.55)          # กรอบภาพประกอบคอลัมน์ขวา
TABLE_BOX = (0.79, 1.32, 8.43, 3.50)
COVER_TITLE_BOX = (1.60, 1.55, 6.80, 1.60)
COVER_HOOK_BOX = (1.60, 3.25, 6.80, 1.10)

# ป้ายชื่อ layout ในเทมเพลต (Google Slides export) + ดัชนีสำรองถ้าไม่พบชื่อ
# key -> (ชื่อ layout ที่ยอมรับได้ เรียงตามลำดับความชอบ, ดัชนีสำรอง)
# เทมเพลตของแต่ละวิชา×ชั้นตั้งชื่อ layout ไม่ตรงกัน จึงต้องมีชื่อสำรอง —
# การถอยไปใช้ "ดัชนี" อย่างเดียวอันตราย เพราะดัชนีเดียวกันคนละเทมเพลตคือคนละ layout
# (เช่น ดัชนี 18: Sci P1 = TITLE_AND_BODY_1 ที่ถูกต้อง แต่ Sci P4 = CUSTOM_3_1 คนละตัว)
LAYOUTS = {
    "title": (("TITLE",), 0),
    "body": (("TITLE_AND_BODY",), 2),
    "title_only": (("TITLE_ONLY",), 4),
    "two_col": (("TITLE_AND_BODY_1", "TITLE_AND_TWO_COLUMNS", "TITLE_AND_BODY"), 18),
    # --- ทางเลือก A: ใช้ layout หลากหลายขึ้นเป็น "พื้นหลัง" (เปิดด้วย --variety)
    # ทุกอันมีพื้นที่ข้อความตำแหน่งเดียวกับ TITLE_AND_BODY ต่างกันแค่ลวดลายมุมจอ
    # จึงสลับได้โดยไม่กระทบการคำนวณกรอบ/การแตกหน้า · เทมเพลตที่ไม่มีชื่อนี้จะถอยไป TITLE_AND_BODY เอง
    "body_v1": (("CUSTOM_18", "TITLE_AND_BODY"), 2),
    "body_v2": (("CUSTOM_18_1", "TITLE_AND_BODY"), 2),
    "body_v3": (("CUSTOM_18_1_2", "TITLE_AND_BODY"), 2),
    "big_point": (("MAIN_POINT", "TITLE_AND_BODY"), 2),
}

ZWSP = dc.ZWSP
_THAI_RE = re.compile(r"[\u0e00-\u0e7f]")


# ================================================================ ฟอนต์: ไฟล์ + เมตริก
@functools.lru_cache(maxsize=None)
def font_path(family: str, bold: bool = False) -> str:
    """หา path ไฟล์ .ttf ของฟอนต์ (ใช้ทั้งวัดขนาดและตรวจสิทธิ์ฝัง)"""
    name = _FONT_FILES.get((family, bold)) or _FONT_FILES.get((family, False))
    if name is None:
        raise KeyError(f"ไม่รู้จักฟอนต์ {family!r}")
    for d in _FONT_DIRS:
        p = os.path.join(d, name)
        if os.path.exists(p):
            return p
    raise FileNotFoundError(
        f"ไม่พบไฟล์ฟอนต์ {name} — ต้องติดตั้ง {family} ก่อน "
        "(ดู ${CLAUDE_PLUGIN_ROOT}/skills/slides/reference/build-notes.md)")


@functools.lru_cache(maxsize=None)
def line_factor(family: str) -> float:
    """ความสูงบรรทัดเดี่ยว (เท่าของขนาดตัวอักษร) จากตาราง hhea/OS/2 ของฟอนต์จริง"""
    with open(font_path(family), "rb") as f:
        d = f.read()
    n = struct.unpack_from(">H", d, 4)[0]
    tabs = {}
    for i in range(n):
        tag, _ck, off, ln = struct.unpack_from(">4sIII", d, 12 + 16 * i)
        tabs[tag.decode("latin1")] = (off, ln)
    upm = struct.unpack_from(">H", d, tabs["head"][0] + 18)[0]
    asc, desc, gap = struct.unpack_from(">hhh", d, tabs["hhea"][0] + 4)
    hhea = (asc - desc + gap) / upm
    win = 0.0
    if "OS/2" in tabs:
        wa, wd = struct.unpack_from(">HH", d, tabs["OS/2"][0] + 74)
        win = (wa + wd) / upm
    return max(hhea, win)


@functools.lru_cache(maxsize=None)
def _cmap(family: str) -> frozenset:
    """เซตของ code point ที่ฟอนต์มี glyph จริง (อ่าน cmap format 4/12)"""
    with open(font_path(family), "rb") as f:
        d = f.read()
    n = struct.unpack_from(">H", d, 4)[0]
    tabs = {}
    for i in range(n):
        tag, _ck, off, ln = struct.unpack_from(">4sIII", d, 12 + 16 * i)
        tabs[tag.decode("latin1")] = (off, ln)
    if "cmap" not in tabs:
        return frozenset()
    base = tabs["cmap"][0]
    ntab = struct.unpack_from(">H", d, base + 2)[0]
    pts, sub4 = set(), None
    for i in range(ntab):
        pid, _eid, so = struct.unpack_from(">HHI", d, base + 4 + 8 * i)
        fmt = struct.unpack_from(">H", d, base + so)[0]
        if fmt == 4 and pid == 3:
            sub4 = base + so
        elif fmt == 12 and pid == 3:
            t = base + so
            ngroups = struct.unpack_from(">I", d, t + 12)[0]
            for g in range(ngroups):
                a, b, _gi = struct.unpack_from(">III", d, t + 16 + 12 * g)
                pts.update(range(a, min(b, a + 0x2000) + 1))
    if sub4 is not None:
        segX2 = struct.unpack_from(">H", d, sub4 + 6)[0]
        for i in range(segX2 // 2):
            end = struct.unpack_from(">H", d, sub4 + 14 + 2 * i)[0]
            start = struct.unpack_from(">H", d, sub4 + 16 + segX2 + 2 * i)[0]
            if start <= end and end != 0xFFFF:
                pts.update(range(start, end + 1))
    return frozenset(pts)


def font_for(ch: str, preferred: str) -> str:
    """ฟอนต์ที่มี glyph ของอักขระนี้จริง

    TH Sarabun New ไม่มีอักษรกรีก (Ω μ π) จึงต้องยืม Prompt ที่ฝังอยู่ในไฟล์แล้ว
    ไม่ปล่อยให้ PowerPoint สลับฟอนต์เอง เพราะฟอนต์ที่ระบบสลับให้จะไม่ถูกฝัง
    -> เปิดเครื่องอื่นแล้วเพี้ยน (และขนาดตัวอักษรก็กระโดด)
    """
    c = ord(ch)
    if c < 0x20 or c == 0x200B or ch.isspace() or c in _cmap(preferred):
        return preferred
    for alt in FONT_FAMILIES:
        if alt != preferred and c in _cmap(alt):
            return alt
    return preferred


def split_by_font(text: str, preferred: str):
    """แตกข้อความเป็นช่วง ๆ ตามฟอนต์ที่รองรับ -> [(ข้อความ, ฟอนต์)]"""
    out = []
    for ch in text:
        fam = font_for(ch, preferred)
        if out and out[-1][1] == fam:
            out[-1][0].append(ch)
        else:
            out.append(([ch], fam))
    return [("".join(buf), fam) for buf, fam in out]


@functools.lru_cache(maxsize=None)
def _pil_font(family: str, bold: bool, px: int):
    from PIL import ImageFont
    return ImageFont.truetype(font_path(family, bold), px)


_MEASURE_SCALE = 8  # โหลดฟอนต์ใหญ่ x8 แล้วหารกลับ เพื่อความละเอียด (PIL รับ size เป็นจำนวนเต็ม)


def text_width_pt(text: str, family: str, size_pt: float, bold: bool = False) -> float:
    """ความกว้างข้อความหน่วย pt (ZWSP ไม่มี glyph ในฟอนต์ ต้องถอดก่อนวัด)"""
    text = text.replace(ZWSP, "")
    if not text:
        return 0.0
    f = _pil_font(family, bold, int(round(size_pt * _MEASURE_SCALE)))
    return f.getlength(text) / _MEASURE_SCALE


# ================================================================ ตัดคำไทย
def display_text(text: str) -> str:
    """ข้อความสุดท้ายที่จะถูกเขียนลงสไลด์จริง (แยกสมการ/ตัวหนา + แทรก ZWSP แล้ว)

    ใช้เป็นฐานของการวัดทุกที่ เพื่อให้ตัวสร้างกับตัวตรวจวัดสตริงเดียวกัน
    """
    out = []
    for kind, span in dc._split_math_spans(text):
        if kind == "math":
            out.append("".join(x[0] for x in _latex_runs(span)))
        else:
            for seg, _bold, _vert in dc._parse_inline(span):
                out.append(thai_break(seg))
    return "".join(out)


def thai_break(text: str) -> str:
    """แทรก ZWSP ตามขอบคำไทย (pythainlp newmm) — PowerPoint จะตัดบรรทัดได้เฉพาะขอบคำจริง"""
    return dc._zwsp_text(text)


_BREAK_RE = re.compile(f"[{ZWSP} ]")


def _chunks(text: str):
    """แตกข้อความเป็นชิ้นที่ 'ห้ามตัดกลาง' พร้อมช่องว่างท้ายชิ้น

    ต้องแตกจาก display_text() ไม่ใช่ข้อความดิบ เพราะ add_text() แยก run ตาม
    $...$ / **หนา** ก่อนแล้วค่อยตัดคำทีละ run — ขอบคำจึงไม่เท่ากับตัดคำทั้งก้อน
    ถ้าวัดคนละสตริงกับที่เขียนจริง การกะหน้าจะคลาด (เคยทำให้ข้อความล้นกรอบ)
    """
    marked = display_text(text)
    out, buf = [], []
    for ch in marked:
        if ch == ZWSP:
            if buf:
                out.append("".join(buf))
                buf = []
            continue
        buf.append(ch)
        if ch == " ":
            out.append("".join(buf))
            buf = []
    if buf:
        out.append("".join(buf))
    return [c for c in out if c]


def wrap_count(text: str, family: str, size_pt: float, width_pt: float,
               bold: bool = False) -> int:
    """จำนวนบรรทัดที่ข้อความจะกินเมื่อจัดในกรอบกว้าง width_pt"""
    if not text:
        return 1
    usable = width_pt * WIDTH_SAFETY
    lines, cur = 1, 0.0
    for ck in _chunks(text):
        w = text_width_pt(ck, family, size_pt, bold)
        if cur > 0 and cur + w > usable:
            lines += 1
            cur = text_width_pt(ck.rstrip(), family, size_pt, bold)
        else:
            cur += w
    return lines


MATH_LINE_BOOST = 1.55  # ย่อหน้าที่มีเศษส่วน/กรณฑ์ ต้องเผื่อความสูงให้สมการซ้อนไม่ชนบรรทัดบน


def has_stacked_math(text: str) -> bool:
    """ย่อหน้านี้มีสมการที่กินความสูง 2 ชั้น (เศษส่วน/กรณฑ์) หรือไม่"""
    return any(kind == "math" and _is_hard_math(span)
               for kind, span in dc._split_math_spans(text or ""))


def line_h_pt(family: str, size_pt: float, boost: float = 1.0) -> float:
    """ความสูงบรรทัด (pt) ที่เราบังคับไว้ — ตรงกับที่ PowerPoint จะจัดจริง"""
    return size_pt * LINE_EXACT[family] * boost


def para_height_pt(text: str, family: str, size_pt: float, width_pt: float,
                   bold: bool = False, space_before: bool = True) -> float:
    """ความสูงที่ย่อหน้าหนึ่งกิน (pt) รวมระยะก่อนย่อหน้า"""
    n = wrap_count(text, family, size_pt, width_pt, bold)
    h = n * line_h_pt(family, size_pt, MATH_LINE_BOOST if has_stacked_math(text) else 1.0)
    if space_before:
        h += size_pt * SPC_BEF_RATIO
    return h


def usable_width_pt(width_in, bullet=True):
    w = (width_in - 2 * TEXT_INSET_IN) * 72
    return w - (BULLET_INDENT_IN * 72 if bullet else 0)


def block_height_pt(items, family, size_pt, width_in, bullet=True):
    """ความสูงรวมของรายการทั้งชุดเมื่อวางในกรอบกว้าง width_in (pt)"""
    width_pt = usable_width_pt(width_in, bullet)
    return sum(para_height_pt(it, family, size_pt, width_pt, space_before=(i > 0))
               for i, it in enumerate(items))


def fits(items, family, size_pt, width_in, height_in, bullet=True) -> bool:
    """ชุดข้อความนี้อยู่ในกรอบได้จริงไหม (ใช้ตัดสินว่าจะวางคู่กับรูปได้หรือไม่)"""
    if not items:
        return True
    limit = (height_in - 2 * TEXT_INSET_IN) * 72 * HEIGHT_SAFETY
    return block_height_pt(items, family, size_pt, width_in, bullet) <= limit


def _greedy_pages(items, family, size_pt, width_pt, limit, cap=None):
    pages, cur, used = [], [], 0.0
    for it in items:
        h = para_height_pt(it, family, size_pt, width_pt, space_before=bool(cur))
        over_cap = cap is not None and len(cur) >= cap
        if cur and (over_cap or used + h > limit):
            pages.append(cur)
            cur, used = [it], para_height_pt(it, family, size_pt, width_pt,
                                             space_before=False)
        else:
            cur.append(it)
            used += h
    if cur:
        pages.append(cur)
    return pages


def fit_pages(items, family, size_pt, width_in, height_in, bullet=True):
    """แบ่ง list ข้อความเป็น 'หน้า' ตามที่วัดจริง — ล้นเมื่อไรขึ้นหน้าใหม่ ไม่ย่อตัวอักษร

    คืน list ของ list (แต่ละอันคือเนื้อหา 1 หน้า) — อย่างน้อยหน้าละ 1 รายการเสมอ
    เมื่อได้หลายหน้าจะลองเกลี่ยให้จำนวนรายการใกล้เคียงกัน (กันหน้าสุดท้ายมีข้อเดียว)
    โดยต้องได้จำนวนหน้าเท่าเดิมเท่านั้น
    """
    if not items:
        return [[]]
    width_pt = (width_in - 2 * TEXT_INSET_IN) * 72
    if bullet:
        width_pt -= BULLET_INDENT_IN * 72
    limit = (height_in - 2 * TEXT_INSET_IN) * 72 * HEIGHT_SAFETY

    pages = _greedy_pages(items, family, size_pt, width_pt, limit)
    if len(pages) > 1:
        cap = -(-len(items) // len(pages))          # เพดานต่อหน้าแบบเกลี่ยเท่า ๆ กัน
        even = _greedy_pages(items, family, size_pt, width_pt, limit, cap=cap)
        if len(even) == len(pages):
            pages = even
    return pages


# ================================================================ เทมเพลต
def template_root() -> str:
    """โฟลเดอร์ Slide Master Template (อยู่ที่รากโปรเจกต์)"""
    here = os.path.dirname(os.path.abspath(__file__))
    root = os.path.abspath(os.path.join(here, "..", "..", "..", ".."))
    return os.path.join(root, "Slide Master Template")


def resolve_template(subject_token: str = "", grade_token: str = "") -> str:
    """หาเทมเพลตตามวิชา×ชั้น: `Slide Master Template/<Subject> <Grade>/Slide Master_<S>_<G>.pptx`

    ลำดับ fallback: ตรงตัว -> วิชาเดียวกันชั้นอื่น -> เทมเพลตแรกที่เจอ (พร้อม WARN)
    """
    root = template_root()
    exact = os.path.join(root, f"{subject_token} {grade_token}",
                         f"Slide Master_{subject_token}_{grade_token}.pptx")
    if os.path.exists(exact):
        return exact
    same_subject = sorted(glob.glob(os.path.join(root, f"{subject_token} *", "*.pptx")))
    if same_subject:
        print(f"WARN: ไม่มีเทมเพลตของ {subject_token} {grade_token} — "
              f"ใช้ {os.path.relpath(same_subject[0], root)} แทน", file=sys.stderr)
        return same_subject[0]
    anyone = sorted(glob.glob(os.path.join(root, "*", "*.pptx")))
    if anyone:
        print(f"WARN: ไม่มีเทมเพลตของ {subject_token} {grade_token} — "
              f"ใช้ {os.path.relpath(anyone[0], root)} แทน", file=sys.stderr)
        return anyone[0]
    raise FileNotFoundError(
        f"ไม่พบเทมเพลตสไลด์ใด ๆ ใน {root} — วางไฟล์ "
        f"'{subject_token} {grade_token}/Slide Master_{subject_token}_{grade_token}.pptx' ก่อน")


def open_template(path: str):
    """เปิดเทมเพลตแล้วลบสไลด์ตัวอย่างทั้งหมด (คง master/layout/ธีมไว้)

    ต้องถอดทั้ง rel และ sldId ไม่งั้นไฟล์สื่อ 14 MB ของสไลด์ตัวอย่างจะติดไปด้วย
    """
    prs = Presentation(path)
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)
    return prs


def get_layout(prs, key: str):
    """หยิบ layout ตามชื่อในเทมเพลต — ลองชื่อที่ยอมรับได้ตามลำดับ แล้วค่อยถอยไปใช้ดัชนี"""
    names, idx = LAYOUTS[key]
    if isinstance(names, str):          # เผื่อสเปกเก่าที่เขียนเป็นสตริงเดี่ยว
        names = (names,)
    master = prs.slide_masters[0]
    by_name = {lay.name: lay for lay in master.slide_layouts}
    avail = [n for n in names if n in by_name]
    if avail:
        # ถ้าวัดความรกไว้แล้ว ให้เลือกอันที่ลวดลายตกแต่งทับพื้นที่ข้อความน้อยที่สุด
        # (เท่ากันให้ยึดลำดับความชอบเดิม) — ไม่มีผลวัดก็ใช้ชื่อแรกที่เจอเหมือนเดิม
        best = avail[0]
        if _CLUTTER:
            scored = [(round(_CLUTTER.get((key, n), 0.0), 3), i, n)
                      for i, n in enumerate(avail)]
            best = min(scored)[2]
        if best != names[0] and (key, best) not in _WARNED:
            _WARNED.add((key, best))          # เตือนครั้งเดียวต่อไฟล์ ไม่ใช่ทุกสไลด์
            print(f"WARN: เลือก layout {best!r} แทน {names[0]!r} "
                  f"(ไม่มีในเทมเพลต หรือลวดลายทับพื้นที่ข้อความมากกว่า)", file=sys.stderr)
        return by_name[best]
    if idx < len(master.slide_layouts):
        print(f"WARN: เทมเพลตไม่มี layout ชื่อใดใน {list(names)} — "
              f"ถอยไปใช้ลำดับที่ {idx} ({master.slide_layouts[idx].name!r}) "
              f"ซึ่งอาจไม่ใช่ layout ที่ต้องการ", file=sys.stderr)
        return master.slide_layouts[idx]
    print(f"WARN: เทมเพลตไม่มีทั้งชื่อใน {list(names)} และลำดับที่ {idx} — ใช้ layout แรก",
          file=sys.stderr)
    return master.slide_layouts[0]


def add_slide(prs, key: str):
    """เพิ่มสไลด์จาก layout แล้ว **ลบ placeholder ที่ติดมา** (เราวางกล่องเองทุกใบ

    เพื่อคุมตำแหน่ง/ขนาด/ฟอนต์ให้เท่ากันทุกหน้า) — กรอบตกแต่งของ layout ยังอยู่ครบ
    """
    slide = prs.slides.add_slide(get_layout(prs, key))
    for shp in list(slide.shapes):
        if shp.is_placeholder:
            shp._element.getparent().remove(shp._element)
    return slide


# ================================================================ ฟอนต์ใน run
def _set_font_tags(rPr, family: str):
    """ตั้ง a:latin + a:ea + a:cs ให้เป็นฟอนต์เดียวกัน (cs = ไทย/complex script)"""
    from lxml import etree
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", family)
    prev = latin
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.Element(qn(tag))
            prev.addnext(el)
        el.set("typeface", family)
        prev = el


def style_run(run, family: str, size_pt: float, bold: bool = False,
              color: RGBColor | None = None, italic: bool = False,
              baseline: int = 0):
    """ตั้งฟอนต์/ขนาด/สี/ยกกำลัง-ตัวห้อย ให้ run เดียว"""
    f = run.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic or None
    f.color.rgb = BODY_INK if color is None else color
    rPr = f._rPr
    _set_font_tags(rPr, family)
    rPr.set("lang", "th-TH")
    rPr.set("altLang", "en-US")
    if baseline:
        rPr.set("baseline", str(baseline))
    return run


# ================================================================ สมการ / สัญลักษณ์
_SUP = 30000   # baseline ของยกกำลัง (30%)
_SUB = -25000  # baseline ของตัวห้อย
DEGREE = "\u00b0"
_OPS_RE = re.compile(r"\s*([\u00d7\u00f7\u00b1=\u2260\u2264\u2265\u2248\u2261])\s*")
_MATH_HARD_RE = re.compile(r"\\(frac|dfrac|tfrac|sqrt|binom|begin|overline|underline|vec|bar|hat)\b")
_FUNC_RE = re.compile(r"^(sin|cos|tan|log|ln|exp|max|min|lim)$")


def _latex_runs(latex: str):
    """LaTeX ง่าย ๆ -> [(text, baseline, italic)] โดยใช้ยกกำลัง/ตัวห้อยจริงของ PowerPoint

    ครอบ x^{2} · H_{2}O · 3.0 \\times 10^{8} · 25\\ \\mathrm{cm^2} · \\mathrm{30^\\circ C}
    (สูตรที่มี \\frac \\sqrt ให้ไปทาง OMML — ดู _is_hard_math)
    """
    s = latex
    s = re.sub(r"\\mathrm\s*\{([^{}]*)\}", lambda m: "\x01" + m.group(1) + "\x02", s)
    s = re.sub(r"\\(?:,|;|:|!|quad|qquad)", " ", s)
    s = s.replace("\\ ", "\x03")                     # \<เว้นวรรค> = ช่องว่างที่สั่งไว้จริง
    s = re.sub(r"(\\[A-Za-z]+) ", r"\1", s)          # LaTeX กลืนช่องว่างที่ปิดท้ายชื่อคำสั่ง
    s = dc._apply_aliases(s).replace("\x03", " ")
    s = s.replace("\u2218", DEGREE)                  # \circ -> องศา ไม่ใช่ ring operator
    s = re.sub(_OPS_RE, r" \1 ", s)                   # คืนช่องว่างรอบเครื่องหมาย = × ÷ ...
    s = re.sub(r" {2,}", " ", s)

    runs = []
    upright = 0
    i = 0
    buf = []

    def flush(base=0):
        if buf:
            txt = "".join(buf)
            for part, ital in _split_italic(txt, upright > 0):
                if part:
                    runs.append((part, base, ital))
            buf.clear()

    while i < len(s):
        c = s[i]
        if c == "\x01":
            flush(); upright += 1; i += 1; continue
        if c == "\x02":
            flush(); upright = max(0, upright - 1); i += 1; continue
        if c in "^_" and i + 1 < len(s):
            j = i + 1
            if s[j] == "{":
                k = s.find("}", j)
                inner, i = (s[j + 1:k], k + 1) if k != -1 else (s[j + 1:], len(s))
            else:
                inner, i = s[j], j + 1
            flush()
            inner = inner.replace("\x01", "").replace("\x02", "").strip()
            if inner == DEGREE:      # ^\circ — องศาลอยอยู่แล้ว ไม่ต้องยกซ้ำ
                runs.append((inner, 0, False))
                continue
            base = _SUP if c == "^" else _SUB
            for part, ital in _split_italic(inner, True):
                if part:
                    runs.append((part, base, ital))
            continue
        if c in "{}":
            i += 1
            continue
        buf.append(c)
        i += 1
    flush()
    return runs


_LETTERS_RE = re.compile(r"[A-Za-z]+")


def _split_italic(text: str, upright: bool):
    """ตัวแปรละตินตัวเดียวให้เอียง (ตามคู่มือ math-symbols-guide) นอกนั้นตัวตรง"""
    if upright:
        return [(text, False)]
    out, pos = [], 0
    for m in _LETTERS_RE.finditer(text):
        if m.start() > pos:
            out.append((text[pos:m.start()], False))
        word = m.group(0)
        out.append((word, len(word) == 1 and not _FUNC_RE.match(word)))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], False))
    return out or [(text, False)]


def _is_hard_math(latex: str) -> bool:
    return bool(_MATH_HARD_RE.search(latex))


def _append_omml(paragraph, latex: str, family: str, size_pt: float):
    """ฝังสมการเป็น OMML จริงของ Office (a14:m) — คืน False ถ้าแปลงไม่ได้"""
    omml = dc.latex_to_omml(latex)
    if omml is None:
        return False
    from lxml import etree
    A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
    MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    math = copy.deepcopy(omml)
    # MML2OMML.XSL ออก w:rPr (ของ Word) — ใน PowerPoint ต้องเป็น a:rPr และต้องระบุขนาด
    # ไม่งั้นสมการจะเล็กกว่าข้อความรอบ ๆ มาก (PowerPoint ใช้ขนาดเริ่มต้นของ placeholder)
    for el in list(math.iter()):
        if el.tag == f"{{{W}}}rPr":
            el.getparent().remove(el)
    for r in math.iter(f"{{{M}}}r"):
        rpr = etree.Element(f"{{{A}}}rPr")
        rpr.set("lang", "th-TH")
        rpr.set("sz", str(int(size_pt * 100)))
        r.insert(0, rpr)
    # ต้องประกาศ prefix a14 ไว้บน mc:Choice เอง เพราะ Requires="a14" อ้างถึง "ชื่อ prefix"
    # (ปล่อยให้ lxml ตั้งชื่อเองเป็น ns1 -> PowerPoint เปิดไฟล์ไม่ได้)
    alt = etree.SubElement(paragraph._p, f"{{{MC}}}AlternateContent", nsmap={"mc": MC})
    choice = etree.SubElement(alt, f"{{{MC}}}Choice", nsmap={"a14": A14})
    choice.set("Requires", "a14")
    a14m = etree.SubElement(choice, f"{{{A14}}}m")
    a14m.append(math)
    fallback = etree.SubElement(alt, f"{{{MC}}}Fallback")
    r = etree.SubElement(fallback, f"{{{A}}}r")
    rPr = etree.SubElement(r, f"{{{A}}}rPr")
    rPr.set("lang", "th-TH")
    rPr.set("sz", str(int(size_pt * 100)))
    _set_font_tags(rPr, family)
    t = etree.SubElement(r, f"{{{A}}}t")
    t.text = "".join(x[0] for x in _latex_runs(latex))
    return True


def put_runs(paragraph, text: str, family: str, size_pt: float, bold: bool = False,
             color: RGBColor | None = None, italic: bool = False, baseline: int = 0):
    """เพิ่ม run พร้อมสลับฟอนต์อัตโนมัติเมื่อฟอนต์หลักไม่มี glyph ของอักขระนั้น"""
    for seg, fam in split_by_font(text, family):
        r = paragraph.add_run()
        r.text = seg
        style_run(r, fam, size_pt, bold, color, italic=italic, baseline=baseline)


def add_text(paragraph, text: str, family: str, size_pt: float, bold: bool = False,
             color: RGBColor | None = None, break_thai: bool = True):
    """เติมข้อความลงย่อหน้า — แยก $...$ เป็นสมการ, `**หนา**`, `^{}`/`_{}` ให้อัตโนมัติ"""
    for kind, span in dc._split_math_spans(text):
        if kind == "math":
            if _is_hard_math(span) and _append_omml(paragraph, span, family, size_pt):
                continue
            for part, base, ital in _latex_runs(span):
                put_runs(paragraph, part, family, size_pt, bold, color,
                         italic=ital, baseline=base)
        else:
            for seg, seg_bold, vert in dc._parse_inline(span):
                base = _SUP if vert == "superscript" else (_SUB if vert == "subscript" else 0)
                put_runs(paragraph, thai_break(seg) if break_thai else seg,
                         family, size_pt, bold or seg_bold, color, baseline=base)
    return paragraph


# ================================================================ กล่องข้อความ
def add_textbox(slide, box_in, anchor=MSO_ANCHOR.TOP):
    """สร้างกล่องข้อความที่ไม่ย่อตัวอักษรเอง (noAutofit) — ล้นเมื่อไรเราขึ้นหน้าใหม่เอง"""
    x, y, w, h = box_in
    shp = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                   Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    bodyPr = tf._txBody.bodyPr
    for attr, val in (("lIns", TEXT_INSET_IN), ("rIns", TEXT_INSET_IN),
                      ("tIns", TEXT_INSET_IN), ("bIns", TEXT_INSET_IN)):
        bodyPr.set(attr, str(int(val * 914400)))
    return shp


def _set_para(par, size_pt: float, bullet: bool, first: bool, align=PP_ALIGN.LEFT,
              family: str = BODY_FONT, boost: float = 1.0):
    par.alignment = align
    par.line_spacing = Pt(line_h_pt(family, size_pt, boost))
    par.space_before = Pt(0 if first else size_pt * SPC_BEF_RATIO)
    par.space_after = Pt(0)
    pPr = par._pPr if par._pPr is not None else par._p.get_or_add_pPr()
    from lxml import etree
    if bullet:
        pPr.set("marL", str(int(BULLET_INDENT_IN * 914400)))
        pPr.set("indent", str(-int(BULLET_INDENT_IN * 914400)))
        buf = etree.SubElement(pPr, qn("a:buFont"))
        buf.set("typeface", "Arial")
        buc = etree.SubElement(pPr, qn("a:buChar"))
        buc.set("char", "\u2022")
    else:
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        etree.SubElement(pPr, qn("a:buNone"))
    return par


_NUMBERED_RE = re.compile(r"^\s*\d+[.)]\s")


def fill_bullets(shape, items, family=BODY_FONT, size_pt=None, bullet=True,
                 color=None, align=PP_ALIGN.LEFT):
    """เทข้อความหลายย่อหน้าลงกล่อง (ย่อหน้าละ 1 บุลเล็ต)"""
    size_pt = size_pt or SIZE["bullet"]
    tf = shape.text_frame
    # ข้อความที่ขึ้นต้นด้วยเลขข้อทุกบรรทัด (ขั้นตอนกิจกรรม/หัวข้อการเรียนรู้) ไม่ต้องมีจุดบุลเล็ตซ้ำ
    if bullet and items and all(_NUMBERED_RE.match(str(x) or "") for x in items):
        bullet = False
    for i, item in enumerate(items):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(par, size_pt, bullet, first=(i == 0), align=align, family=family,
                  boost=MATH_LINE_BOOST if has_stacked_math(item) else 1.0)
        add_text(par, item, family, size_pt, color=color)
    return shape


def title_metrics(text: str, wide: bool = False):
    """คำนวณขนาด/กรอบของหัวข้อหน้า โดยไม่สร้างรูปทรง

    ลดขนาดลงทีละขั้นให้พอ 1 บรรทัดก่อน ถ้าไม่ไหวจึงยอม 2 บรรทัด (แล้วดันเนื้อหาลง)
    คืน (size_pt, box, content_top_in)
    """
    box = list(TITLE_BOX_WIDE if wide else TITLE_BOX)
    width_pt = (box[2] - 2 * TEXT_INSET_IN) * 72
    size = SIZE["title"]
    while size > SIZE["title_min"] and wrap_count(text, HEAD_FONT, size, width_pt, True) > 1:
        size -= 2
    lines = wrap_count(text, HEAD_FONT, size, width_pt, True)
    box[3] = lines * line_h_pt(HEAD_FONT, size) / 72 + 2 * TEXT_INSET_IN
    return size, box, max(BODY_BOX[1], box[1] + box[3] + 0.10)


def add_title(slide, text: str, wide: bool = False):
    """วางหัวข้อหน้า (Prompt Bold) — คืน y (นิ้ว) ที่เนื้อหาควรเริ่ม"""
    size, box, top = title_metrics(text, wide)
    shp = add_textbox(slide, box)
    par = shp.text_frame.paragraphs[0]
    _set_para(par, size, bullet=False, first=True, family=HEAD_FONT)
    add_text(par, text, HEAD_FONT, size, bold=True, color=INK)
    return top


def set_notes(slide, *parts):
    """บันทึกโน้ตผู้สอน (speaker note + คำสั่งวาดภาพ) ลงหน้าโน้ต"""
    text = "\n\n".join(p for p in parts if p)
    if not text:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for par in tf.paragraphs:
        for run in par.runs:
            style_run(run, BODY_FONT, 12)


# ================================================================ กรอบภาพ / รูป
def add_image_frame(slide, box_in, caption="ภาพประกอบ", image_file=None):
    """มีไฟล์รูป -> ฝังรูปพอดีกรอบ ; ไม่มี -> กรอบมนจาง ๆ ไว้ให้เติมทีหลัง"""
    from pptx.enum.shapes import MSO_SHAPE
    x, y, w, h = box_in
    if image_file and os.path.exists(image_file):
        from PIL import Image
        with Image.open(image_file) as im:
            iw, ih = im.size
        scale = min(w / (iw / 96), h / (ih / 96))
        dw, dh = (iw / 96) * scale, (ih / 96) * scale
        return slide.shapes.add_picture(
            image_file, Emu(int((x + (w - dw) / 2) * 914400)),
            Emu(int((y + (h - dh) / 2) * 914400)),
            Emu(int(dw * 914400)), Emu(int(dh * 914400)))

    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x * 914400)), Emu(int(y * 914400)),
        Emu(int(w * 914400)), Emu(int(h * 914400)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BAND
    shp.line.color.rgb = SOFT
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    par = tf.paragraphs[0]
    _set_para(par, SIZE["caption"], bullet=False, first=True, align=PP_ALIGN.CENTER)
    add_text(par, caption, BODY_FONT, SIZE["caption"], color=FRAME_MUTED)
    return shp


# ================================================================ ตาราง
def auto_col_ratio(rows, size_pt=None):
    """สัดส่วนความกว้างคอลัมน์ตามความยาวข้อความจริง — หัวคอลัมน์สั้น ๆ จะไม่ตกบรรทัด"""
    size_pt = size_pt or SIZE["table_cell"]
    pad = 2 * 0.08 * 72
    widths = []
    for c in range(len(rows[0])):
        w = 0.0
        for r, row in enumerate(rows):
            head = (r == 0)
            w = max(w, text_width_pt(row[c], HEAD_FONT if head else BODY_FONT,
                                     SIZE["table_head"] if head else size_pt, head))
        widths.append(w + pad)
    total = sum(widths) or 1.0
    return [w / total for w in widths]


def add_table(slide, rows, box_in, col_ratio=None, size_pt=None):
    """ตารางคำศัพท์ — แถวแรกเป็นหัวตาราง (Prompt) ที่เหลือ TH Sarabun New"""
    size_pt = size_pt or SIZE["table_cell"]
    x, y, w, h = box_in
    n_rows, n_cols = len(rows), len(rows[0])
    gf = slide.shapes.add_table(n_rows, n_cols, Emu(int(x * 914400)),
                                Emu(int(y * 914400)), Emu(int(w * 914400)),
                                Emu(int(h * 914400)))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    ratio = col_ratio or auto_col_ratio(rows, size_pt)
    for c, frac in enumerate(ratio):
        tbl.columns[c].width = Emu(int(w * frac * 914400))
    row_h = h / n_rows
    for r, row in enumerate(rows):
        tbl.rows[r].height = Emu(int(row_h * 914400))
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Emu(int(0.08 * 914400))
            cell.margin_top = cell.margin_bottom = Emu(int(CELL_PAD_IN / 2 * 914400))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND if r == 0 else (BAND if r % 2 == 0 else WHITE)
            par = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            _set_para(par, size_pt, bullet=False, first=True,
                      align=PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT,
                      family=HEAD_FONT if r == 0 else BODY_FONT)
            add_text(par, val, HEAD_FONT if r == 0 else BODY_FONT,
                     SIZE["table_head"] if r == 0 else size_pt,
                     bold=(r == 0), color=WHITE if r == 0 else TABLE_INK)
    return tbl


CELL_PAD_IN = 0.06  # ขอบบน+ล่างของเซลล์ตาราง


def table_row_h_in(size_pt=None) -> float:
    """ความสูงแถวตาราง (นิ้ว) — บรรทัดเดียว + ขอบเซลล์"""
    size_pt = size_pt or SIZE["table_cell"]
    return line_h_pt(BODY_FONT, size_pt) / 72 + CELL_PAD_IN


def table_rows_fit(rows, box_h_in, size_pt=None) -> int:
    """จำนวนแถวสูงสุดที่ใส่ได้ในกรอบ (รวมแถวหัวตาราง)"""
    return max(2, int((box_h_in * HEIGHT_SAFETY) // table_row_h_in(size_pt)))


# ================================================================ จานสีตามพื้นเทมเพลต
# ทำไมต้องมี: สีตัวอักษรถูกจูนไว้กับเทมเพลตพื้นสว่าง (Sci P1) พอผู้ใช้วางเทมเพลต
# พื้นกระดานดำ (Sci P4) ข้อความสีเข้มก็จมหายไปกับพื้นทั้งแผ่น — อ่านไม่ได้เลยแม้แต่ตัวเดียว
# จึงต้องให้ตัวสร้าง "ดูพื้นก่อนเลือกสีหมึก" แทนการ hard-code ไว้ชุดเดียว
_PALETTES = {
    "light": {
        "INK": RGBColor(0x0E, 0x2A, 0x47),   # น้ำเงินเข้ม
        "BODY_INK": RGBColor(0x20, 0x20, 0x20),
        "SOFT": RGBColor(0xD5, 0xD9, 0xEE),
        "MUTED": RGBColor(0x86, 0x9F, 0xB2),
    },
    "dark": {
        "INK": RGBColor(0xFF, 0xFF, 0xFF),   # ขาว — หัวข้อบนพื้นเข้ม
        "BODY_INK": RGBColor(0xF2, 0xF2, 0xF2),
        "SOFT": RGBColor(0x9A, 0xA6, 0xC8),
        "MUTED": RGBColor(0xC9, 0xD6, 0xDF),
    },
}
_PALETTE_NOW = "light"


def apply_palette(kind: str) -> None:
    """สลับจานสีหมึกทั้งโมดูล ('light' | 'dark') — ต้องเรียกก่อนวาดสไลด์"""
    global INK, BODY_INK, SOFT, MUTED, _PALETTE_NOW
    p = _PALETTES[kind]
    INK, BODY_INK, SOFT, MUTED = p["INK"], p["BODY_INK"], p["SOFT"], p["MUTED"]
    _PALETTE_NOW = kind


def _cache_path() -> str:
    root = os.path.dirname(os.path.abspath(__file__))
    for _ in range(4):
        root = os.path.dirname(root)
    d = os.path.join(root, ".cache")
    os.makedirs(d, exist_ok=True)
    return os.path.join(d, "pptx_bg.json")


def detect_bg_dark(template_path: str) -> bool:
    """พื้นของเทมเพลตนี้เข้มหรือไม่ — วัดจากภาพจริง ไม่เดาจาก XML

    พื้นหลังของเทมเพลตมาได้หลายทาง (ธีม · รูปเต็มแผ่นบน master · รูปทรงตกแต่งซ้อนกัน)
    การไล่อ่าน XML จึงเปราะมาก วิธีที่เชื่อได้คือ **render สไลด์เปล่าออกมาดูพิกเซลจริง**
    แล้วจำผลไว้ในแคช (คีย์ = path + mtime) เพื่อไม่ต้องเปิด PowerPoint ซ้ำทุกครั้ง
    """
    key = "%s|%d" % (os.path.abspath(template_path),
                     int(os.path.getmtime(template_path)))
    cache = {}
    cp = _cache_path()
    if os.path.exists(cp):
        try:
            with open(cp, encoding="utf-8") as f:
                cache = json.load(f)
        except Exception:
            cache = {}
    if key in cache:
        return bool(cache[key])

    dark = False
    tmp = None
    try:
        import tempfile
        from PIL import Image
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        import embed_fonts_pptx as ef

        tmp = tempfile.mkdtemp(prefix="bgprobe_")
        probe = os.path.join(tmp, "probe.pptx")
        prs = open_template(template_path)
        add_slide(prs, "body")          # สไลด์เปล่าที่มีแต่กรอบตกแต่งของ layout
        prs.save(probe)
        if ef.export_png(probe, tmp, width=640) != 0:
            raise RuntimeError("export PNG ไม่สำเร็จ")
        png = os.path.join(tmp, "slide01.png")
        im = Image.open(png).convert("L")
        w, h = im.size
        # เก็บเฉพาะบริเวณที่ข้อความไปวางจริง (เว้นขอบที่มักเป็นลวดลายตกแต่ง)
        box = im.crop((int(w * 0.12), int(h * 0.15), int(w * 0.88), int(h * 0.85)))
        px = sorted(box.getdata())
        med = px[len(px) // 2]
        dark = med < 128
        print("พื้นเทมเพลต: ความสว่างกลางภาพ %d/255 -> %s"
              % (med, "เข้ม (ใช้หมึกสว่าง)" if dark else "สว่าง (ใช้หมึกเข้ม)"))
    except Exception as exc:
        print("WARN: ตรวจความสว่างพื้นเทมเพลตไม่ได้ (%s) — ใช้จานสีพื้นสว่างตามค่าเริ่มต้น"
              % exc, file=sys.stderr)
        return False
    finally:
        if tmp:
            import shutil
            shutil.rmtree(tmp, ignore_errors=True)

    cache[key] = dark
    try:
        with open(cp, "w", encoding="utf-8") as f:
            json.dump(cache, f, ensure_ascii=False, indent=1)
    except Exception:
        pass
    return dark


# ================================================================ เลือก layout ที่พื้นที่ข้อความโล่งที่สุด
# ทำไมต้องวัดจากภาพ: เทมเพลตแต่ละเล่มวางลวดลายตกแต่งไม่เหมือนกัน บางเล่ม (Sci P4)
# วางรูปดินสอทับตรงคอลัมน์ข้อความพอดี ส่วนบางเล่ม (Sci P1) ใช้กรอบขอบเต็มแผ่น
# ซึ่ง "กรอบสี่เหลี่ยม" ของมันกินทั้งจอแต่ตรงกลางโล่ง — ตัดสินจาก bounding box จึงผิดทั้งคู่
# ต้อง render สไลด์เปล่าของแต่ละ layout แล้วนับพิกเซลที่ต่างจากพื้น เฉพาะบริเวณที่ข้อความจะไปวาง
# หมายเหตุ: ไม่นับ IMAGE_BOX เพราะเราวางการ์ดทึบ/รูปจริงทับตรงนั้นอยู่แล้ว
# ลวดลายที่อยู่หลังการ์ดจึงมองไม่เห็น ไม่ควรทำให้ layout นั้นเสียคะแนน
_KEY_BOXES = {
    "title": (COVER_TITLE_BOX, COVER_HOOK_BOX),
    "body": (TITLE_BOX, BODY_BOX),
    "title_only": (TITLE_BOX_WIDE,),
    "two_col": (TITLE_BOX_WIDE, COL_L_BOX),
    "body_v1": (TITLE_BOX, BODY_BOX),
    "body_v2": (TITLE_BOX, BODY_BOX),
    "body_v3": (TITLE_BOX, BODY_BOX),
    "big_point": (TITLE_BOX, BODY_BOX),
}
_CLUTTER_VER = 3        # เปลี่ยนเลขนี้เมื่อแก้นิยามกล่องหรือเพิ่ม key เพื่อทิ้งแคชเก่า
SLIDE_W_IN, SLIDE_H_IN = 10.0, 5.625


def _clutter_score(img, boxes):
    """สัดส่วนพิกเซลในบริเวณข้อความที่ 'ไม่ใช่พื้น' (= มีลวดลายตกแต่งมาทับ)"""
    w, h = img.size
    allpx = sorted(img.getdata())
    base = allpx[len(allpx) // 2]           # ความสว่างพื้นของสไลด์นี้
    worst = 0.0
    for (x, y, bw, bh) in boxes:
        c = img.crop((int(w * x / SLIDE_W_IN), int(h * y / SLIDE_H_IN),
                      int(w * (x + bw) / SLIDE_W_IN), int(h * (y + bh) / SLIDE_H_IN)))
        px = list(c.getdata())
        if not px:
            continue
        off = sum(1 for v in px if abs(v - base) > 40) / len(px)
        worst = max(worst, off)             # กล่องที่แย่ที่สุดเป็นตัวตัดสิน
    return worst


def _layout_clutter(template_path: str) -> dict:
    """คืน {(key, ชื่อ layout): คะแนนความรก} — วัดครั้งเดียวต่อเทมเพลตแล้วแคชไว้"""
    key = "clutter%d|%s|%d" % (_CLUTTER_VER, os.path.abspath(template_path),
                               int(os.path.getmtime(template_path)))
    cp = _cache_path()
    cache = {}
    if os.path.exists(cp):
        try:
            with open(cp, encoding="utf-8") as f:
                cache = json.load(f)
        except Exception:
            cache = {}
    if key in cache:
        return {tuple(k.split("\t")): v for k, v in cache[key].items()}

    out = {}
    tmp = None
    try:
        import tempfile
        from PIL import Image
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        import embed_fonts_pptx as ef

        prs = open_template(template_path)
        by_name = {L.name: L for L in prs.slide_masters[0].slide_layouts}
        order = []
        for k, (names, _idx) in LAYOUTS.items():
            for nm in (names if not isinstance(names, str) else (names,)):
                if nm in by_name:
                    prs.slides.add_slide(by_name[nm])
                    order.append((k, nm))
        if not order:
            return {}
        for sl in prs.slides:                # เอา placeholder ออกให้เหลือแต่ลวดลาย
            for shp in list(sl.shapes):
                if shp.is_placeholder:
                    shp._element.getparent().remove(shp._element)
        tmp = tempfile.mkdtemp(prefix="clutter_")
        probe = os.path.join(tmp, "probe.pptx")
        prs.save(probe)
        if ef.export_png(probe, tmp, width=800) != 0:
            raise RuntimeError("export PNG ไม่สำเร็จ")
        for i, (k, nm) in enumerate(order, 1):
            png = os.path.join(tmp, "slide%02d.png" % i)
            out[(k, nm)] = _clutter_score(Image.open(png).convert("L"), _KEY_BOXES[k])
    except Exception as exc:
        print("WARN: วัดความรกของ layout ไม่ได้ (%s) — ใช้ลำดับชื่อตามที่ตั้งไว้" % exc,
              file=sys.stderr)
        return {}
    finally:
        if tmp:
            import shutil
            shutil.rmtree(tmp, ignore_errors=True)

    cache[key] = {"\t".join(k): v for k, v in out.items()}
    try:
        with open(cp, "w", encoding="utf-8") as f:
            json.dump(cache, f, ensure_ascii=False, indent=1)
    except Exception:
        pass
    return out


_CLUTTER: dict = {}
_WARNED: set = set()


def use_template(template_path: str) -> None:
    """เตรียมทุกอย่างที่ขึ้นกับเทมเพลต: จานสีหมึก + คะแนนความรกของแต่ละ layout"""
    global _CLUTTER
    _WARNED.clear()
    apply_palette("dark" if detect_bg_dark(template_path) else "light")
    _CLUTTER = _layout_clutter(template_path)
    for (k, nm), v in sorted(_CLUTTER.items()):
        if v > 0.02:
            print("  layout %-24s (%s) พื้นที่ข้อความมีลวดลายทับ %.0f%%" % (nm, k, v * 100))
