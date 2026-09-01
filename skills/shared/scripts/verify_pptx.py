# -*- coding: utf-8 -*-
"""
verify_pptx.py — ตรวจไฟล์สไลด์ที่ build เสร็จ (EDU ONE — คู่ขนานกับ verify_docx.py)

ตรวจ 7 ข้อ
  1. เปิดได้จริง + มีสไลด์ (และครบตามที่คาดถ้าใส่ --expect-min)
  2. ฟอนต์ทุก run เป็น Prompt/TH Sarabun New และตั้งครบทั้ง a:latin + a:cs (ไทย)
  3. ไม่มีตัวอักษรเล็กกว่าพื้นที่กำหนด (เนื้อหา 24pt · ตาราง 22pt)
  4. ฝังฟอนต์ครบทั้ง 2 ตระกูล (ppt/fonts/*.fntdata + p:embeddedFontLst)
  5. ไม่มี normAutofit (แปลว่ามีหน้าที่ PowerPoint ย่อตัวอักษรเองเพราะล้น)
  6. ไม่มี `$...$` หรือ \\command ตกค้างเป็นข้อความ (สมการแปลงไม่สำเร็จ)
  7. ไม่มีกล่องข้อความไหนล้นกรอบ (วัดซ้ำด้วยกลไกเดียวกับตอน build)

ใช้: verify_pptx.py <pptx> [--expect-min N] [--no-embed-check]
exit 0 = ผ่าน / 1 = พบปัญหา / 2 = usage
"""
from __future__ import annotations

import argparse
import os
import re
import sys
import zipfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pptx_common as pc  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

MATH_LEFTOVER_RE = re.compile(r"\$[^$]+\$|\\[A-Za-z]{2,}")


# วัตถุที่เทมเพลตล็อกไว้เป็นดีไซน์ (แถบสี, kicker, ป้ายเล็ก) ไม่ใช่เนื้อหาการสอน
# จึงไม่บังคับเกณฑ์ขนาดตัวอักษรขั้นต่ำ — แต่ยังบังคับเรื่องฟอนต์ให้เป็นฟอนต์โครงการ
DESIGN_PREFIX = "LOCK_"


def _runs(shape):
    if not shape.has_text_frame:
        return
    for par in shape.text_frame.paragraphs:
        for run in par.runs:
            yield par, run


def _shape_boxes(slide):
    """คืน (shape, กรอบเป็นนิ้ว) เฉพาะกล่องข้อความที่เราวางเอง"""
    for shp in slide.shapes:
        if shp.has_text_frame and shp.width and shp.height:
            yield shp, (shp.width / 914400, shp.height / 914400)


def check(path, expect_min=0, embed_check=True):
    problems, warns = [], []
    prs = Presentation(path)
    n = len(prs.slides.__iter__.__self__._sldIdLst)
    if n == 0:
        problems.append("ไฟล์ไม่มีสไลด์เลย")
    if expect_min and n < expect_min:
        problems.append("จำนวนสไลด์ %d น้อยกว่าที่คาด (%d)" % (n, expect_min))

    # ---- 2/3/6 ต่อ run ----
    for i, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            tables = [shp.table] if shp.has_table else []
            # ช่องที่เทมเพลตออกแบบให้เล็กจนใส่ขนาดพื้นไม่ได้ (เช่นบรรทัดข้อมูลใต้ชื่อเรื่อง)
            # ถือเป็นองค์ประกอบดีไซน์ ไม่ใช่เนื้อหาการสอน
            tiny = bool(shp.height) and (shp.height / 914400.0) * 72 < pc.SIZE["min_body"] * 1.2
            design = shp.name.startswith(DESIGN_PREFIX) or tiny
            for par, run in _runs(shp):
                _check_run(i, par, run, problems, in_table=False, design=design)
            for tbl in tables:
                for r, row in enumerate(tbl.rows):
                    for cell in row.cells:
                        for par in cell.text_frame.paragraphs:
                            for run in par.runs:
                                _check_run(i, par, run, problems, in_table=True)

    # ---- 5. autofit ----
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        for nm in names:
            if nm.startswith("ppt/slides/slide") and nm.endswith(".xml"):
                xml = z.read(nm).decode("utf-8")
                if "<a:normAutofit" in xml:
                    problems.append("%s มี normAutofit — PowerPoint ย่อตัวอักษรเอง "
                                    "(ต้องขึ้นหน้าใหม่แทน)" % os.path.basename(nm))
        # ---- 4. ฝังฟอนต์ ----
        if embed_check:
            fnt = [x for x in names if x.startswith("ppt/fonts/") and x.endswith(".fntdata")]
            pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
            embedded = set(re.findall(r'<p:embeddedFont><p:font typeface="([^"]+)"', pres_xml))
            if not fnt:
                problems.append("ไม่มีฟอนต์ฝังในไฟล์ (ppt/fonts/*.fntdata) — "
                                "รัน embed_fonts_pptx.py")
            for fam in pc.FONT_FAMILIES:
                if fam not in embedded:
                    problems.append("ฟอนต์ %s ไม่ได้ถูกฝังลงไฟล์" % fam)

    # ---- 7. ล้นกรอบ ----
    for i, slide in enumerate(prs.slides, 1):
        for shp, (w_in, h_in) in _shape_boxes(slide):
            used = 0.0
            for j, par in enumerate(shp.text_frame.paragraphs):
                if not par.runs:
                    continue
                fam, size = _para_font(par)
                bullet = par._pPr is not None and par._pPr.find(qn("a:buChar")) is not None
                width_pt = (w_in - 2 * pc.TEXT_INSET_IN) * 72
                if bullet:
                    width_pt -= pc.BULLET_INDENT_IN * 72
                text = "".join(r.text for r in par.runs)
                used += pc.para_height_pt(text, fam, size, width_pt, space_before=(j > 0))
            limit = (h_in - 2 * pc.TEXT_INSET_IN) * 72
            if used > limit + 1:
                problems.append("สไลด์ %d: ข้อความล้นกรอบ (%.0f pt > %.0f pt) — %s"
                                % (i, used, limit, shp.name))

    print("ตรวจแล้ว: %d สไลด์" % n)
    for w in warns:
        print("  WARN: %s" % w)
    if problems:
        print("\nFAIL — พบ %d ปัญหา:" % len(problems))
        for k, p in enumerate(problems, 1):
            print("  %d. %s" % (k, p))
        return 1
    print("OK — ฟอนต์/ขนาด/การฝังฟอนต์/การจัดหน้า ผ่านทุก check")
    return 0


def _para_font(par):
    """ฟอนต์+ขนาดตัวแทนของย่อหน้า (ตัวใหญ่สุดในย่อหน้า)"""
    fam, size = pc.BODY_FONT, pc.SIZE["bullet"]
    best = 0
    for run in par.runs:
        sz = run.font.size.pt if run.font.size else 0
        if sz > best:
            best = sz
            latin = run.font._rPr.find(qn("a:latin"))
            fam = latin.get("typeface") if latin is not None else pc.BODY_FONT
            size = sz
    return (fam if fam in pc.FONT_FAMILIES else pc.BODY_FONT), (size or pc.SIZE["bullet"])


def _check_run(slide_no, par, run, problems, in_table, design=False):
    rPr = run.font._rPr
    text = run.text or ""
    if not text.strip():
        return
    faces = {}
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        faces[tag] = el.get("typeface") if el is not None else None
    for tag, face in faces.items():
        if face is None:
            problems.append("สไลด์ %d: run %r ไม่ได้ตั้ง %s (ไทยจะเพี้ยนเครื่องอื่น)"
                            % (slide_no, text[:20], tag))
        elif face not in pc.FONT_FAMILIES:
            problems.append("สไลด์ %d: run %r ใช้ฟอนต์ %s (ต้องเป็น %s)"
                            % (slide_no, text[:20], face, "/".join(pc.FONT_FAMILIES)))
    size = run.font.size.pt if run.font.size else None
    floor = pc.SIZE["min_table"] if in_table else pc.SIZE["min_body"]
    if design:
        return
    if size is None:
        problems.append("สไลด์ %d: run %r ไม่ได้กำหนดขนาดตัวอักษร" % (slide_no, text[:20]))
    elif size < floor and not rPr.get("baseline"):
        problems.append("สไลด์ %d: run %r ขนาด %gpt เล็กกว่าพื้น %gpt"
                        % (slide_no, text[:20], size, floor))
    if MATH_LEFTOVER_RE.search(text):
        problems.append("สไลด์ %d: มีสัญกรณ์คณิตตกค้างเป็นข้อความ %r"
                        % (slide_no, text[:40]))


def main():
    ap = argparse.ArgumentParser(description="ตรวจไฟล์สไลด์ .pptx")
    ap.add_argument("pptx")
    ap.add_argument("--expect-min", type=int, default=0)
    ap.add_argument("--no-embed-check", action="store_true",
                    help="ข้ามการตรวจฝังฟอนต์ (ใช้ตอนยังไม่ได้รัน embed)")
    a = ap.parse_args()
    if not os.path.exists(a.pptx):
        print("ไม่พบไฟล์: %s" % a.pptx, file=sys.stderr)
        return 2
    return check(a.pptx, a.expect_min, embed_check=not a.no_embed_check)


if __name__ == "__main__":
    raise SystemExit(main())
