# -*- coding: utf-8 -*-
"""
pptx_slots.py — engine สำหรับ "เทมเพลตแบบสเปก LAYOUT" (EDU ONE — agent 4)

เทมเพลตแบบนี้ (เช่น `Mat M1`) ไม่ได้ให้มาเป็น slideLayout เปล่า ๆ แต่ให้มาเป็น
**สไลด์ตัวอย่างหนึ่งหน้าต่อหนึ่ง LAYOUT** ซึ่งวางตำแหน่ง/สี/ฟอนต์ไว้เรียบร้อย และตั้งชื่อวัตถุ
ที่แก้ได้ไว้ชัดเจน (`TITLE_*`, `BODY_*`, `IMAGE_*`, `TABLE_*`, `FORMULA_*`, `META_*` ฯลฯ)
พร้อมกฎกำกับในหน้า เช่น "แก้ไขเฉพาะวัตถุชื่อ ... | ห้ามยืดภาพ"

วิธีทำงานของ engine นี้จึงเป็น **โคลนหน้าตัวอย่างของ LAYOUT ที่ต้องการ แล้วเติมข้อความลงวัตถุตามชื่อ**
ไม่ใช่วาดกล่องเองที่พิกัดตายตัวแบบ engine เดิม (ซึ่งใช้กับเทมเพลตที่ดีไซน์อยู่ใน slideLayout)

`build_slides.py` เลือก engine อัตโนมัติจาก `is_slot_template()`

รันผ่านตัวห่อ: eduone-py <ชื่อไฟล์นี้> <args>  (หา Python 3.12 ให้เองทุก OS)
"""
from __future__ import annotations

import copy
import os
import re
import sys

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pptx_common as pc  # noqa: E402

# ---------------------------------------------------------------- คำใบ้ในเทมเพลต
LAYOUT_CODE_SHAPE = "LOCK_LAYOUT_CODE"      # กล่องที่เขียนว่า "LAYOUT 01"
_CODE_RE = re.compile(r"LAYOUT\s*0*(\d+)", re.IGNORECASE)

# วัตถุที่เป็น "คู่มือสำหรับผู้สร้างสื่อ" ต้องไม่ติดไปกับสไลด์ที่นักเรียนเห็น
GUIDE_SHAPES = {
    "LAYOUT_NAME", "LOCK_LAYOUT_CODE", "LOCK_USAGE_NOTE",
    "LOCK_FOOTER", "LOCK_SAFE_AREA_NOTE",
}

# ---------------------------------------------------------------- ขนาดตัวอักษร
# เทมเพลตตั้งข้อความตัวอย่างไว้ 16pt บนสไลด์ 13.33 นิ้ว (= 12pt บนสไลด์ 10 นิ้ว) ซึ่งเล็กเกินไป
# สำหรับการฉายหน้าห้อง — ผู้ใช้เลือก "ยึดกฎอ่านจากหลังห้อง" จึงคงกรอบ/สี/ตำแหน่งของเทมเพลตไว้
# แล้วขยายเฉพาะขนาดตัวอักษร (อ้างอิงสไลด์กว้าง 13.333 นิ้ว — สเกลตามความกว้างจริงของเทมเพลต)
REF_SLIDE_W_IN = 13.333
SIZE = {
    "cover_title": 44,
    "cover_hook": 30,
    "cover_meta": 22,
    "micro": 12,          # บรรทัดข้อมูลเล็กที่เทมเพลตออกแบบให้เป็น caption
    "title": 40,          # หัวข้อหน้า — ลดลงได้ถึง title_min เพื่อให้พอ 1 บรรทัด
    "title_min": 30,
    "body": 30,           # เนื้อหาหลักในกล่องใหญ่
    "body_min": 24,
    "card": 24,           # กล่องเล็ก (การ์ด/ขั้นตอน/สถานการณ์) — พื้นที่จำกัด
    "card_min": 24,       # พื้นของการ์ด = เกณฑ์เดียวกับ verify — ยาวเกินให้ถอย layout
    "card_title": 26,
    "caption": 24,
    "meta": 24,
    "table_head": 24,
    "table_cell": 24,
    "formula": 34,        # สมการเด่นกลางหน้า
}

_BULLET_RE = re.compile(r"^\s*[•\-•]\s*")


def scale_for(prs) -> float:
    """สเกลขนาดตัวอักษรตามความกว้างจริงของเทมเพลต (อ้างอิง 13.333 นิ้ว)"""
    return (prs.slide_width / 914400.0) / REF_SLIDE_W_IN


# ================================================================ ตรวจชนิดเทมเพลต
def layout_index(prs) -> dict:
    """คืน {เลข LAYOUT: slide ตัวอย่าง} จากสไลด์ในเทมเพลต"""
    out = {}
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.name != LAYOUT_CODE_SHAPE or not shp.has_text_frame:
                continue
            m = _CODE_RE.search(shp.text_frame.text or "")
            if m:
                out[int(m.group(1))] = slide
            break
    return out


def is_slot_template(prs) -> bool:
    """เทมเพลตนี้เป็นแบบสเปก LAYOUT (มีหน้าตัวอย่างพร้อมรหัส LAYOUT) หรือไม่"""
    return len(layout_index(prs)) >= 3


# ================================================================ โคลน/ล้างหน้า
def clone(prs, src):
    """โคลนสไลด์ตัวอย่างมาเป็นหน้าใหม่ท้ายเด็ค (ยกทุกวัตถุมาทั้งดุ้น)"""
    dst = prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes):          # ลบ placeholder ที่ add_slide แถมมา
        shp._element.getparent().remove(shp._element)
    tree = dst.shapes._spTree
    for shp in src.shapes:
        if shp.shape_type is not None and shp.has_text_frame is False and shp.shape_type == 13:
            print("WARN: เทมเพลตมีรูปฝังในหน้าตัวอย่าง — engine นี้ยังไม่ย้าย rel ให้",
                  file=sys.stderr)
        tree.append(copy.deepcopy(shp._element))
    strip_guides(dst)
    normalize_fonts(dst)
    return dst


def normalize_fonts(slide) -> None:
    """บังคับฟอนต์ของข้อความที่ติดมากับเทมเพลตให้เป็นฟอนต์โครงการ

    เทมเพลตใช้ Calibri ตามธีม Office — ถ้าปล่อยไว้ ข้อความตกแต่ง (เช่น kicker บนหน้าปก)
    จะเป็นคนละฟอนต์กับเนื้อหา และภาษาไทยจะเพี้ยนเมื่อเปิดเครื่องอื่น
    ช่องที่เราเติมข้อความเองจะถูกตั้งฟอนต์ทับอีกทีอยู่แล้ว
    """
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        fam = pc.HEAD_FONT if shp.name.startswith("TITLE") else pc.BODY_FONT
        for par in shp.text_frame.paragraphs:
            for run in par.runs:
                pc._set_font_tags(run.font._rPr, fam)
                run.font._rPr.set("lang", "th-TH")


def strip_guides(slide) -> None:
    """เอาข้อความคู่มือของเทมเพลตออก (LAYOUT 01 / กฎการใช้ / footer)"""
    for shp in list(slide.shapes):
        if shp.name in GUIDE_SHAPES:
            shp._element.getparent().remove(shp._element)


def drop_samples(prs, keep_from: int) -> None:
    """ลบสไลด์ตัวอย่างของเทมเพลต (ต้องเรียกหลังโคลนครบทุกหน้าแล้ว)"""
    lst = prs.slides._sldIdLst
    for sld_id in list(lst)[:keep_from]:
        prs.part.drop_rel(sld_id.rId)
        lst.remove(sld_id)


def slots(slide) -> dict:
    """{ชื่อวัตถุ: shape} ของหน้านี้"""
    return {shp.name: shp for shp in slide.shapes}


def drop(shape) -> None:
    """เอาวัตถุออกจากหน้า — เรียกซ้ำกับวัตถุเดิมได้ (ไม่ระเบิด)"""
    if shape is None:
        return
    parent = shape._element.getparent()
    if parent is not None:
        parent.remove(shape._element)


def drop_unused(slide, used) -> None:
    """ลบช่องที่ไม่ได้เติม เพื่อไม่ให้ข้อความตัวอย่าง [ ... ] หลุดไปถึงนักเรียน"""
    for name, shp in list(slots(slide).items()):
        if name in used or not shp.has_text_frame:
            continue
        t = (shp.text_frame.text or "").strip()
        if t.startswith("[") or not t:      # ข้อความตัวอย่างของเทมเพลต
            drop(shp)


# ================================================================ เติมข้อความ
def box_in(shape):
    """กรอบของวัตถุเป็นนิ้ว (x, y, w, h)"""
    return (shape.left / 914400.0, shape.top / 914400.0,
            shape.width / 914400.0, shape.height / 914400.0)


def fit_size(text, family, start, minimum, width_in, height_in, lines_max=None,
             bullet=False):
    """เลือกขนาดตัวอักษรที่ใหญ่ที่สุดที่ยังอยู่ในกรอบ (ไม่ต่ำกว่า minimum)

    รับได้ทั้ง str และ list[str] — ถ้าเป็นรายการต้องวัดแบบหลายย่อหน้า (มีระยะก่อนย่อหน้า)
    ให้ตรงกับที่ verify_pptx วัด ไม่งั้นจะคำนวณต่างกันแล้วข้อความล้นกรอบ
    """
    items = [text] if isinstance(text, str) else list(text)
    usable_h = (height_in - 2 * pc.TEXT_INSET_IN) * 72
    size = start
    while size > minimum:
        h = pc.block_height_pt(items, family, size, width_in, bullet=bullet)
        n = max(pc.wrap_count(x, family, size,
                              pc.usable_width_pt(width_in, bullet)) for x in items)
        if (lines_max is None or n <= lines_max) and h <= usable_h:
            return size
        size -= 2
    return minimum


def _fill_of(shape):
    """คัดลอก <a:solidFill> ของ run แรกในช่อง — คือ **สีที่เทมเพลตออกแบบไว้**

    ห้ามบังคับสีเอง เพราะแต่ละช่องมีพื้นหลังต่างกัน (เช่น หัวข้อบนแถบสีเข้มต้องเป็นตัวอักษรสีอ่อน)
    """
    if not shape.has_text_frame:
        return None
    for par in shape.text_frame.paragraphs:
        for run in par.runs:
            f = run.font._rPr.find(qn("a:solidFill"))
            if f is not None:
                return copy.deepcopy(f)
        pPr = par._pPr
        if pPr is not None:
            d = pPr.find(qn("a:defRPr"))
            if d is not None:
                f = d.find(qn("a:solidFill"))
                if f is not None:
                    return copy.deepcopy(f)
    return None


def set_text(shape, items, size, family=None, bold=False, color=None,
             align=PP_ALIGN.LEFT, bullet=False, anchor=None, scale=1.0):
    """ล้างข้อความตัวอย่างแล้วเติมของจริง — ใช้กลไกฟอนต์/ตัดคำ/สมการชุดเดียวกับ engine เดิม

    **สีตัวอักษรสืบทอดจากเทมเพลตเสมอ** (ยกเว้นผู้เรียกระบุ color มาเอง)
    """
    if isinstance(items, str):
        items = [items]
    items = [_BULLET_RE.sub("", str(x)) for x in items if str(x).strip()]
    if not items:
        drop(shape)
        return None
    family = family or pc.BODY_FONT
    size = max(1, round(size * scale))
    keep_fill = _fill_of(shape) if color is None else None
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if anchor is not None:
        tf.vertical_anchor = anchor
    bodyPr = tf._txBody.bodyPr
    for attr in ("lIns", "rIns", "tIns", "bIns"):
        bodyPr.set(attr, str(int(pc.TEXT_INSET_IN * 914400)))
    for i, item in enumerate(items):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pc._set_para(par, size, bullet, first=(i == 0), align=align, family=family,
                     boost=pc.MATH_LINE_BOOST if pc.has_stacked_math(item) else 1.0)
        pc.add_text(par, item, family, size, bold=bold, color=color)
        if keep_fill is not None:
            for run in par.runs:                      # คืนสีของเทมเพลตให้ทุก run
                rPr = run.font._rPr
                old_fill = rPr.find(qn("a:solidFill"))
                if old_fill is not None:
                    rPr.remove(old_fill)
                latin = rPr.find(qn("a:latin"))
                fill = copy.deepcopy(keep_fill)
                if latin is not None:
                    latin.addprevious(fill)
                else:
                    rPr.append(fill)
    return shape


def put_image(slide, shape, image_path, scale=1.0):
    """วางรูปให้พอดีกรอบของช่องภาพ **แบบ Fit ไม่ยืด ไม่ crop** ตามกฎของเทมเพลต

    ไม่มีไฟล์ -> คงกรอบเดิมไว้เป็นที่ว่างพร้อมป้าย (ไม่ทำให้ build พัง)
    """
    x, y, w, h = box_in(shape)
    if not image_path or not os.path.exists(image_path):
        set_text(shape, "ภาพประกอบ", SIZE["caption"], color=pc.MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, scale=scale)
        return None
    from PIL import Image
    with Image.open(image_path) as im:
        iw, ih = im.size
    ratio = min(w / iw, h / ih)
    dw, dh = iw * ratio, ih * ratio
    pic = slide.shapes.add_picture(
        image_path, Emu(int((x + (w - dw) / 2) * 914400)),
        Emu(int((y + (h - dh) / 2) * 914400)),
        Emu(int(dw * 914400)), Emu(int(dh * 914400)))
    drop(shape)
    return pic


def fill_table(slide, shape, rows, scale=1.0):
    """เติมตารางที่เทมเพลตวางไว้แล้ว (ปรับจำนวนแถวให้พอดีข้อมูล)"""
    tbl = shape.table
    need = len(rows)
    # python-pptx ไม่รองรับ index ติดลบของ tbl.rows — ต้องอ้างด้วยตำแหน่งจริง
    while len(tbl.rows) > need:                       # ลบแถวเกิน
        tr = tbl.rows[len(tbl.rows) - 1]._tr
        tr.getparent().remove(tr)
    while len(tbl.rows) < need:                       # เพิ่มแถวโดยโคลนแถวสุดท้าย
        tr = copy.deepcopy(tbl.rows[len(tbl.rows) - 1]._tr)
        tbl._tbl.append(tr)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c >= len(tbl.columns):
                break
            cell = tbl.cell(r, c)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            head = (r == 0)
            size = round((SIZE["table_head"] if head else SIZE["table_cell"]) * scale)
            par = tf.paragraphs[0]
            pc._set_para(par, size, bullet=False, first=True,
                         align=PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT,
                         family=pc.HEAD_FONT if head else pc.BODY_FONT)
            pc.add_text(par, str(val), pc.HEAD_FONT if head else pc.BODY_FONT, size,
                        bold=head, color=pc.WHITE if head else pc.BODY_INK)
    return tbl


def set_notes(slide, *parts):
    pc.set_notes(slide, *parts)


# ================================================================ ข้อจำกัดของเทมเพลต
_SIZE_OF_SLOT = (
    ("TITLE", "title"), ("HOOK", "cover_hook"), ("FORMULA_MAIN", "formula"),
    ("BODY", "body"), ("TOPIC_LIST", "body"), ("SUMMARY_LIST", "body"),
    ("EXPLANATION", "card"), ("PROMPT", "body"), ("ANSWER", "body"),
    ("META", "meta"), ("CAPTION", "caption"),
)


def _size_key_for(name):
    for token, key in _SIZE_OF_SLOT:
        if token in name:
            return key
    return "card"


def limits(prs):
    """ตารางข้อจำกัดของทุก LAYOUT — ชื่อช่อง · ขนาดกรอบ · จำนวนตัวอักษรไทยที่ใส่ได้

    ใช้ส่งให้ slides-writer อ่านก่อนเขียน เพื่อไม่ให้เขียนยาวเกินกรอบที่ดีไซน์ไว้
    (พิสูจน์แล้วว่าลดการถอย layout จาก 6 หน้าเหลือ 0)
    """
    scale = scale_for(prs)
    out = []
    for code, slide in sorted(layout_index(prs).items()):
        name = ""
        for shp in slide.shapes:
            if shp.name == "LAYOUT_NAME" and shp.has_text_frame:
                name = shp.text_frame.text.strip()
                break
        rows = []
        for shp in slide.shapes:
            if shp.name.startswith(("LOCK_", "LAYOUT_")) or not shp.has_text_frame:
                continue
            x, y, w, h = box_in(shp)
            key = _size_key_for(shp.name)
            size = max(1, round(SIZE[key] * scale))
            per_line = int(pc.usable_width_pt(w, bullet=False)
                           / max(1.0, pc.text_width_pt("ก", pc.BODY_FONT, size)))
            lines = max(1, int(((h - 2 * pc.TEXT_INSET_IN) * 72) // pc.line_h_pt(pc.BODY_FONT, size)))
            rows.append((shp.name, w, h, size, per_line, lines, per_line * lines))
        out.append((code, name, rows))
    return out


def print_limits(template_path):
    prs = Presentation(template_path)
    if not is_slot_template(prs):
        print("เทมเพลตนี้ไม่ใช่แบบสเปก LAYOUT — ไม่มีข้อจำกัดรายช่องให้แสดง")
        return 0
    print("ข้อจำกัดของเทมเพลต: %s" % os.path.basename(template_path))
    print("(จำนวนตัวอักษรคิดที่ขนาดพื้นของช่องนั้น ภาษาไทยความกว้างเฉลี่ย)")
    for code, name, rows in limits(prs):
        print("\nLAYOUT %02d — %s" % (code, name))
        for nm, w, h, size, per_line, lines, total in rows:
            print("   %-24s กรอบ %4.2f x %4.2f นิ้ว · %2dpt · ~%3d ตัว/บรรทัด · %d บรรทัด · รวม ~%d ตัว"
                  % (nm, w, h, size, per_line, lines, total))
    return 0


def main(argv):
    if len(argv) == 3 and argv[1] == "--limits":
        return print_limits(argv[2])
    print(__doc__, file=sys.stderr)
    return 2


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
